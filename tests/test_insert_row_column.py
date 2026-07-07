"""insert_row / insert_column (v0.14) — 위치 지정 삽입 + 서식 상속 + 병합 가드.

외부 리소스 없이 pytest 로 돌아간다:
    pytest tests/test_insert_row_column.py -v
"""
from __future__ import annotations

from pathlib import Path

import pytest
from docx import Document
from docx.oxml.ns import qn
from hwpx.document import HwpxDocument
from pptx import Presentation
from pptx.util import Inches

from document_adapter import load
from document_adapter.base import NotImplementedForFormat

FILL = "D9E2F3"


# -------- 픽스처 --------

def _shade(cell, fill: str = FILL) -> None:
    """DOCX 셀에 배경 음영(w:shd) 지정 — 서식 상속 검증용."""
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.makeelement(qn("w:shd"), {qn("w:val"): "clear", qn("w:fill"): fill})
    tc_pr.append(shd)


def _make_docx_year_table(path: Path) -> None:
    """사진8 형태 축소판: [구분/2025/2024/2023] × 2열, 라벨 열 음영."""
    doc = Document()
    table = doc.add_table(rows=4, cols=2)
    for r, label in enumerate(["구분", "2025", "2024", "2023"]):
        table.cell(r, 0).text = label
        table.cell(r, 1).text = f"v{r}"
        _shade(table.cell(r, 0))
    doc.save(path)


def _make_hwpx_table(path: Path) -> None:
    """3x3 HWPX 표."""
    doc = HwpxDocument.new()
    doc.add_paragraph("")
    doc.add_table(3, 3)
    doc.save_to_path(path)

    doc2 = HwpxDocument.open(path)
    try:
        tbl = next(t for p in doc2.sections[0].paragraphs for t in p.tables)
        for r in range(3):
            for c in range(3):
                tbl.cell(r, c).text = f"r{r}c{c}"
        doc2.save_to_path(path)
    finally:
        doc2.close()


def _make_pptx_table(path: Path, rows: int = 3, cols: int = 3) -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(
        rows, cols, Inches(1), Inches(1), Inches(8), Inches(3)
    )
    for r in range(rows):
        for c in range(cols):
            shape.table.cell(r, c).text = f"r{r}c{c}"
    prs.save(path)


# -------- DOCX --------

def test_docx_insert_row_middle_with_format_inheritance(tmp_path: Path) -> None:
    """2025 위(인덱스 1)에 2026 삽입 — 위치·값·음영 상속을 모두 검증."""
    src = tmp_path / "years.docx"
    _make_docx_year_table(src)

    adapter = load(src)
    try:
        idx = adapter.insert_row(0, ["2026", "new"], at_row=1)
        adapter.save()
    finally:
        adapter.close()
    assert idx == 1

    verify = load(src)
    try:
        t = verify.get_tables(preview_rows=10)[0]
    finally:
        verify.close()
    assert t.rows == 5
    assert t.preview[1] == ["2026", "new"]
    assert t.preview[2] == ["2025", "v1"]  # 기존 행은 아래로 밀림
    assert t.preview[4] == ["2023", "v3"]

    # 서식 상속: 새 행 라벨 셀에 템플릿 행의 음영이 복사됐는지
    doc = Document(str(src))
    new_tc = doc.tables[0].rows[1].cells[0]._tc
    shd = new_tc.tcPr.find(qn("w:shd"))
    assert shd is not None
    assert shd.get(qn("w:fill")) == FILL


def test_docx_append_row_inherits_formatting(tmp_path: Path) -> None:
    """append_row(=insert_row 위임)도 마지막 행 서식을 상속해야 한다 (v0.14 수정)."""
    src = tmp_path / "append.docx"
    _make_docx_year_table(src)

    adapter = load(src)
    try:
        adapter.append_row(0, ["2022", "old"])
        adapter.save()
    finally:
        adapter.close()

    verify = load(src)
    try:
        t = verify.get_tables(preview_rows=10)[0]
    finally:
        verify.close()
    assert t.rows == 5
    assert t.preview[4] == ["2022", "old"]

    doc = Document(str(src))
    new_tc = doc.tables[0].rows[4].cells[0]._tc
    shd = new_tc.tcPr.find(qn("w:shd"))
    assert shd is not None and shd.get(qn("w:fill")) == FILL


def test_docx_insert_row_rejects_vertical_merge_crossing(tmp_path: Path) -> None:
    src = tmp_path / "vmerge.docx"
    doc = Document()
    table = doc.add_table(rows=3, cols=2)
    table.cell(0, 0).merge(table.cell(1, 0))  # (0,0)~(1,0) 세로 병합
    doc.save(src)

    adapter = load(src)
    try:
        with pytest.raises(NotImplementedForFormat):
            adapter.insert_row(0, ["x", "y"], at_row=1)
    finally:
        adapter.close()


def test_docx_insert_column_middle_rebalances_width(tmp_path: Path) -> None:
    src = tmp_path / "cols.docx"
    doc = Document()
    table = doc.add_table(rows=3, cols=3)
    for r in range(3):
        for c in range(3):
            table.cell(r, c).text = f"r{r}c{c}"
    doc.save(src)

    before = Document(str(src))
    grid = before.tables[0]._tbl.tblGrid
    total_before = sum(int(gc.get(qn("w:w"))) for gc in grid.gridCol_lst)

    adapter = load(src)
    try:
        idx = adapter.insert_column(0, ["신규", "a", "b"], at_col=1)
        adapter.save()
    finally:
        adapter.close()
    assert idx == 1

    verify = load(src)
    try:
        t = verify.get_tables(preview_rows=10)[0]
    finally:
        verify.close()
    assert t.cols == 4
    assert t.preview[0] == ["r0c0", "신규", "r0c1", "r0c2"]
    assert t.preview[1] == ["r1c0", "a", "r1c1", "r1c2"]
    assert t.preview[2] == ["r2c0", "b", "r2c1", "r2c2"]

    # 표 전체 폭 유지 (반올림 오차 허용)
    after = Document(str(src))
    grid2 = after.tables[0]._tbl.tblGrid
    cols_after = grid2.gridCol_lst
    assert len(cols_after) == 4
    total_after = sum(int(gc.get(qn("w:w"))) for gc in cols_after)
    assert abs(total_after - total_before) <= len(cols_after)


def test_docx_insert_column_rejects_horizontal_merge_crossing(
    tmp_path: Path,
) -> None:
    src = tmp_path / "hmerge.docx"
    doc = Document()
    table = doc.add_table(rows=2, cols=3)
    a = table.cell(0, 0)
    a.merge(table.cell(0, 1))  # (0,0)~(0,1) 가로 병합 — 경계 1 을 가로지름
    doc.save(src)

    adapter = load(src)
    try:
        with pytest.raises(NotImplementedForFormat):
            adapter.insert_column(0, ["x", "y"], at_col=1)
    finally:
        adapter.close()


# -------- HWPX --------

def test_hwpx_insert_row_middle(tmp_path: Path) -> None:
    src = tmp_path / "insert.hwpx"
    _make_hwpx_table(src)

    adapter = load(src)
    try:
        idx = adapter.insert_row(0, ["N0", "N1", "N2"], at_row=1)
        adapter.save()
    finally:
        adapter.close()
    assert idx == 1

    verify = load(src)
    try:
        t = verify.get_tables(preview_rows=10)[0]
    finally:
        verify.close()
    assert t.rows == 4
    assert t.preview[0] == ["r0c0", "r0c1", "r0c2"]
    assert t.preview[1] == ["N0", "N1", "N2"]
    assert t.preview[2] == ["r1c0", "r1c1", "r1c2"]
    assert t.preview[3] == ["r2c0", "r2c1", "r2c2"]


def test_hwpx_insert_column_middle(tmp_path: Path) -> None:
    src = tmp_path / "insertcol.hwpx"
    _make_hwpx_table(src)

    adapter = load(src)
    try:
        idx = adapter.insert_column(0, ["H", "a", "b"], at_col=1)
        adapter.save()
    finally:
        adapter.close()
    assert idx == 1

    verify = load(src)
    try:
        t = verify.get_tables(preview_rows=10)[0]
    finally:
        verify.close()
    assert t.cols == 4
    assert t.preview[0] == ["r0c0", "H", "r0c1", "r0c2"]
    assert t.preview[1] == ["r1c0", "a", "r1c1", "r1c2"]
    assert t.preview[2] == ["r2c0", "b", "r2c1", "r2c2"]


def test_hwpx_insert_column_at_end(tmp_path: Path) -> None:
    src = tmp_path / "appendcol.hwpx"
    _make_hwpx_table(src)

    adapter = load(src)
    try:
        idx = adapter.insert_column(0, ["x", "y", "z"])
        adapter.save()
    finally:
        adapter.close()
    assert idx == 3

    verify = load(src)
    try:
        t = verify.get_tables(preview_rows=10)[0]
    finally:
        verify.close()
    assert t.cols == 4
    assert t.preview[0] == ["r0c0", "r0c1", "r0c2", "x"]
    assert t.preview[2] == ["r2c0", "r2c1", "r2c2", "z"]


# -------- PPTX --------

def test_pptx_insert_row_middle(tmp_path: Path) -> None:
    src = tmp_path / "insert.pptx"
    _make_pptx_table(src)

    adapter = load(src)
    try:
        idx = adapter.insert_row(0, ["N0", "N1", "N2"], at_row=1)
        adapter.save()
    finally:
        adapter.close()
    assert idx == 1

    verify = load(src)
    try:
        t = verify.get_tables(preview_rows=10)[0]
    finally:
        verify.close()
    assert t.rows == 4
    assert t.preview[1] == ["N0", "N1", "N2"]
    assert t.preview[2] == ["r1c0", "r1c1", "r1c2"]


def test_pptx_insert_column_middle(tmp_path: Path) -> None:
    src = tmp_path / "insertcol.pptx"
    _make_pptx_table(src)

    adapter = load(src)
    try:
        idx = adapter.insert_column(0, ["H", "a", "b"], at_col=1)
        adapter.save()
    finally:
        adapter.close()
    assert idx == 1

    verify = load(src)
    try:
        t = verify.get_tables(preview_rows=10)[0]
    finally:
        verify.close()
    assert t.cols == 4
    assert t.preview[0] == ["r0c0", "H", "r0c1", "r0c2"]
    assert t.preview[1] == ["r1c0", "a", "r1c1", "r1c2"]


# -------- 미지원 포맷 기본값 --------

def test_xlsx_insert_defaults_raise(tmp_path: Path) -> None:
    """base 기본 구현은 NotImplementedForFormat — XLSX 등 미구현 포맷 보호."""
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    src = tmp_path / "sheet.xlsx"
    wb = Workbook()
    wb.active["A1"] = "x"
    wb.save(src)

    adapter = load(src)
    try:
        if hasattr(type(adapter), "insert_row") and (
            type(adapter).insert_row is not None
        ):
            # XlsxAdapter 가 자체 구현을 갖기 전까지 base 기본값이 적용된다
            from document_adapter.base import DocumentAdapter

            if type(adapter).insert_row is DocumentAdapter.insert_row:
                with pytest.raises(NotImplementedForFormat):
                    adapter.insert_row(0, ["a"], at_row=0)
    finally:
        adapter.close()
