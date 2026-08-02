"""PPTX shape 단위 복사 (v0.18): copy_shape."""
from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from document_adapter import load
from document_adapter.base import NotImplementedForFormat
from document_adapter.tools import call_tool


def _make_deck(path: Path) -> None:
    """3장 덱: 표지(placeholder) + 표 슬라이드 + 차트 슬라이드."""
    prs = Presentation()

    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "표지 제목"
    s1.placeholders[1].text = "부제목"

    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "표 페이지"
    tbl = s2.shapes.add_table(
        3, 2, Inches(1), Inches(2), Inches(6), Inches(2)
    ).table
    tbl.cell(0, 0).text = "항목"
    tbl.cell(0, 1).text = "내용"
    tbl.cell(1, 0).text = "담당자"
    tbl.cell(1, 1).text = "유지수"
    tbl.cell(2, 0).text = "상태"
    tbl.cell(2, 1).text = "진행중"

    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    s3.shapes.title.text = "차트 페이지"
    cd = CategoryChartData()
    cd.categories = ["1분기", "2분기"]
    cd.add_series("매출", (100.0, 200.0))
    s3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(2), Inches(8), Inches(4), cd,
    )
    prs.save(path)


@pytest.fixture()
def deck(tmp_path: Path) -> Path:
    p = tmp_path / "deck.pptx"
    _make_deck(p)
    return p


# ---- 표 복사 ----

def test_copy_table_clear_values(deck: Path) -> None:
    doc = load(deck)
    try:
        result = doc.copy_shape(3, table_index=0, clear_values=True)
        doc.save()
    finally:
        doc.close()
    assert result["kind"] == "table"
    assert result["values_cleared"] is True
    assert result["source_slide_index"] == 2
    # 새 표의 preview 는 전부 빈 값 (서식/구조만 유지)
    assert all(not (c or "").strip()
               for row in result["preview"] for c in row if c is not None)

    doc = load(deck)
    try:
        tables = doc.get_tables()
        assert len(tables) == 2
        # 원본(slide 2) 내용 불변
        orig = next(t for t in tables if t.location == "slide 2")
        assert orig.preview[1][1] == "유지수"
        # 복사본(slide 3)은 3x2 구조 유지 + 빈 값 — 반환 좌표로 바로 편집 가능
        new_tidx = result["table_index"]
        assert doc.get_cell(new_tidx, 1, 0).text == ""
        doc.set_cell(new_tidx, 1, 0, "새담당")
        assert doc.get_cell(new_tidx, 1, 0).text == "새담당"
    finally:
        doc.close()


def test_copy_table_keep_values(deck: Path) -> None:
    doc = load(deck)
    try:
        result = doc.copy_shape(1, table_index=0)
        doc.save()
    finally:
        doc.close()
    assert result["values_cleared"] is False
    # slide 1 에 복사됐으므로 전역 순번상 새 표가 0번 → 시프트 warning
    assert result["table_index"] == 0
    assert "warning" in result

    doc = load(deck)
    try:
        assert doc.get_cell(0, 1, 1).text == "유지수"   # 복사본 (slide 1)
        assert doc.get_cell(1, 1, 1).text == "유지수"   # 원본 (slide 2)
    finally:
        doc.close()


def test_copy_table_style_preserved(deck: Path) -> None:
    """복사본 표의 XML 구조(스타일 참조 포함)가 원본과 동질이어야 한다."""
    doc = load(deck)
    try:
        doc.copy_shape(3, table_index=0, clear_values=True)
        doc.save()
    finally:
        doc.close()
    prs = Presentation(str(deck))

    def tbl_style(slide):
        for sh in slide.shapes:
            if getattr(sh, "has_table", False):
                el = sh.table._tbl
                pr = el.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}tblPr")
                return None if pr is None else pr.findtext(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId")
        return None

    assert tbl_style(prs.slides[1]) == tbl_style(prs.slides[2])


# ---- 차트 복사 ----

def test_copy_chart_independent(deck: Path) -> None:
    doc = load(deck)
    try:
        src_chart = doc.get_charts(slide_index=3)[0]
        result = doc.copy_shape(
            2, source_slide_index=3, shape_id=src_chart.shape_id,
        )
        # 복사본만 수치 편집 → 원본 불변이어야 함
        doc.set_chart_data(
            2, result["shape_id"],
            set_points=[{"series": "매출", "category": "1분기", "value": 999}],
        )
        doc.save()
    finally:
        doc.close()
    assert result["kind"] == "chart"
    assert result["chart_type"] == "COLUMN_CLUSTERED"
    assert "set_chart_data" in result["hint"]

    doc = load(deck)
    try:
        copied = doc.get_charts(slide_index=2)[0]
        orig = doc.get_charts(slide_index=3)[0]
        assert copied.series[0]["values"] == [999.0, 200.0]
        assert orig.series[0]["values"] == [100.0, 200.0]
    finally:
        doc.close()

    # chart part / 내장 워크북까지 분리됐는지
    prs = Presentation(str(deck))
    parts = set()
    for sl in prs.slides:
        for sh in sl.shapes:
            if getattr(sh, "has_chart", False):
                parts.add(str(sh.chart.part.partname))
    assert len(parts) == 2


# ---- 텍스트박스 / placeholder ----

def test_copy_placeholder_becomes_regular_shape(deck: Path) -> None:
    """placeholder 복사 → <p:ph> 제거 + 위치/크기 실측 고정."""
    doc = load(deck)
    try:
        title = doc.get_shapes(slide_index=2)[0]  # "표 페이지" 제목 placeholder
        result = doc.copy_shape(
            3, source_slide_index=2, shape_id=title.shape_id,
            x_cm=1.0, y_cm=6.0,
        )
        doc.save()
    finally:
        doc.close()
    assert result["kind"] == "text"

    prs = Presentation(str(deck))
    new_shape = next(
        sh for sh in prs.slides[2].shapes if sh.shape_id == result["shape_id"]
    )
    assert new_shape.is_placeholder is False
    assert new_shape.text_frame.text == "표 페이지"
    assert abs(new_shape.left - 360000) < 1000     # 1.0 cm
    assert abs(new_shape.top - 2160000) < 1000     # 6.0 cm
    assert new_shape.width is not None


def test_copy_text_clear_values(deck: Path) -> None:
    doc = load(deck)
    try:
        title = doc.get_shapes(slide_index=1)[0]
        result = doc.copy_shape(
            2, source_slide_index=1, shape_id=title.shape_id,
            clear_values=True, y_cm=15.0,
        )
        doc.save()
        old = doc.set_shape_text(2, result["shape_id"], "새 텍스트")
        assert old == ""   # 값이 비워져 있었음
        doc.save()
    finally:
        doc.close()


# ---- id 유일성 / 검증 ----

def test_copy_twice_unique_ids(deck: Path) -> None:
    doc = load(deck)
    try:
        r1 = doc.copy_shape(3, table_index=0, clear_values=True)
        r2 = doc.copy_shape(3, table_index=0, clear_values=True, y_cm=12.0)
        doc.save()
    finally:
        doc.close()
    assert r1["shape_id"] != r2["shape_id"]

    prs = Presentation(str(deck))
    ids = [sh.shape_id for sh in prs.slides[2].shapes]
    assert len(ids) == len(set(ids))


def test_copy_invalid_args(deck: Path) -> None:
    doc = load(deck)
    try:
        with pytest.raises(ValueError):
            doc.copy_shape(2)                          # 원본 미지정
        with pytest.raises(ValueError):
            doc.copy_shape(2, table_index=0,
                           source_slide_index=3, shape_id=4)  # 동시 지정
        with pytest.raises(ValueError):
            doc.copy_shape(99, table_index=0)          # target 범위 밖
        with pytest.raises(ValueError):
            doc.copy_shape(2, source_slide_index=1, shape_id=9999)
        with pytest.raises(IndexError):
            doc.copy_shape(2, table_index=99)
    finally:
        doc.close()


def test_non_pptx_rejected(tmp_path: Path) -> None:
    from docx import Document

    p = tmp_path / "doc.docx"
    Document().save(str(p))
    doc = load(p)
    try:
        with pytest.raises(NotImplementedForFormat):
            doc.copy_shape(1, table_index=0)
    finally:
        doc.close()


def test_tool_dispatch_copy_shape(deck: Path) -> None:
    result = call_tool("copy_shape", {
        "path": str(deck), "target_slide_index": 3,
        "table_index": 0, "clear_values": True,
    })
    assert "error" not in result
    assert result["kind"] == "table"

    edited = call_tool("set_cell", {
        "path": str(deck), "table_index": result["table_index"],
        "row": 1, "col": 1, "value": "홍길동",
    })
    assert "error" not in edited
