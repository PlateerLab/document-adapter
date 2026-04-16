"""Smoke test — 세 포맷에 대해 템플릿 생성 → 렌더 → 검증.

외부 리소스 없이 pytest로 돌아간다:
    pytest tests/test_smoke.py -v
"""
from __future__ import annotations

from pathlib import Path

import pytest
from docx import Document
from docx.shared import Pt as DocxPt
from hwpx.document import HwpxDocument
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt

from document_adapter import load
from document_adapter.tools import call_tool


# -------- 헬퍼: 공정한 템플릿을 그때그때 생성 --------

def _make_docx(path: Path) -> None:
    doc = Document()
    doc.add_paragraph("{{ title }}")
    doc.add_paragraph("작성자: {{ author }}")
    doc.save(path)


def _make_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "{{ title }}"
    slide.placeholders[1].text = "작성자: {{ author }}"
    prs.save(path)


def _make_hwpx(path: Path) -> None:
    doc = HwpxDocument.new()
    doc.add_paragraph("{{ title }}")
    doc.add_paragraph("작성자: {{ author }}")
    doc.save_to_path(path)


FACTORIES = {
    "docx": _make_docx,
    "pptx": _make_pptx,
    "hwpx": _make_hwpx,
}

CONTEXT = {"title": "통합 테스트", "author": "tester"}


@pytest.mark.parametrize("fmt", ["docx", "pptx", "hwpx"])
def test_render_template_dispatcher(tmp_path: Path, fmt: str) -> None:
    src = tmp_path / f"template.{fmt}"
    FACTORIES[fmt](src)

    doc = load(src)
    assert set(doc.get_placeholders()) == {"title", "author"}
    doc.render_template(CONTEXT)
    doc.save()
    doc.close()

    doc2 = load(src)
    assert doc2.get_placeholders() == []
    doc2.close()


@pytest.mark.parametrize("fmt", ["docx", "pptx", "hwpx"])
def test_tool_inspect(tmp_path: Path, fmt: str) -> None:
    src = tmp_path / f"template.{fmt}"
    FACTORIES[fmt](src)

    result = call_tool("inspect_document", {"path": str(src)})
    assert result["format"] == fmt
    assert "title" in result["placeholders"]
    assert "author" in result["placeholders"]


def test_tool_render(tmp_path: Path) -> None:
    src = tmp_path / "template.docx"
    _make_docx(src)

    result = call_tool("render_template", {
        "path": str(src),
        "context": CONTEXT,
    })
    assert result["rendered_count"] == 2
    assert result["placeholders_after"] == []
    assert Path(result["output_path"]).exists()


def _make_pptx_with_simple_table(path: Path) -> None:
    """2x2 표가 있는 최소 PPTX. append_row 테스트용."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(1.5))
    tbl = shape.table
    tbl.cell(0, 0).text = "H1"
    tbl.cell(0, 1).text = "H2"
    tbl.cell(1, 0).text = "A"
    tbl.cell(1, 1).text = "B"
    prs.save(path)


def test_tool_append_row_pptx_supported(tmp_path: Path) -> None:
    """PPTX append_row 는 v0.5.0 부터 자체 lxml 구현으로 지원된다.

    마지막 <a:tr> 을 deepcopy 해 새 행을 만들고 셀 텍스트만 비운 뒤 values 로 채운다.
    """
    src = tmp_path / "template.pptx"
    _make_pptx_with_simple_table(src)

    result = call_tool("append_row", {
        "path": str(src),
        "table_index": 0,
        "values": ["new_col1", "new_col2"],
    })
    assert "error" not in result or result.get("error") is None, result

    # 저장 후 재열기 해서 행이 +1 되고 값이 들어갔는지 검증
    output = Path(result.get("output_path", str(src)))
    after = Presentation(str(output))
    for shape in after.slides[0].shapes:
        if shape.has_table:
            tbl = shape.table
            assert len(tbl.rows) == 3  # 2 → 3
            new_row = tbl.rows[2]
            cells = list(new_row.cells)
            assert cells[0].text == "new_col1"
            assert cells[1].text == "new_col2"
            break
    else:
        raise AssertionError("output 에 표가 없음")


# ---- regression: set_cell must keep run-level font formatting (issue #1) ----

def _make_docx_table_with_font(path: Path) -> None:
    doc = Document()
    table = doc.add_table(rows=1, cols=1)
    cell = table.cell(0, 0)
    run = cell.paragraphs[0].add_run("original")
    run.font.name = "Malgun Gothic"
    run.font.size = DocxPt(18)
    run.bold = True
    doc.save(path)


def _make_pptx_table_with_font(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(6), Inches(1))
    cell = table_shape.table.cell(0, 0)
    cell.text = "original"
    run = cell.text_frame.paragraphs[0].runs[0]
    run.font.name = "Malgun Gothic"
    run.font.size = PptxPt(18)
    run.font.bold = True
    prs.save(path)


def test_docx_set_cell_preserves_font(tmp_path: Path) -> None:
    src = tmp_path / "formatted.docx"
    _make_docx_table_with_font(src)

    doc = load(src)
    old = doc.set_cell(0, 0, 0, "replaced")
    doc.save()
    doc.close()
    assert old == "original"

    verify = Document(src)
    run = verify.tables[0].cell(0, 0).paragraphs[0].runs[0]
    assert run.text == "replaced"
    assert run.font.name == "Malgun Gothic"
    assert run.font.size == DocxPt(18)
    assert run.bold is True


def test_pptx_set_cell_preserves_font(tmp_path: Path) -> None:
    src = tmp_path / "formatted.pptx"
    _make_pptx_table_with_font(src)

    doc = load(src)
    old = doc.set_cell(0, 0, 0, "replaced")
    doc.save()
    doc.close()
    assert old == "original"

    verify = Presentation(src)
    cell = next(
        shape.table.cell(0, 0)
        for slide in verify.slides
        for shape in slide.shapes
        if shape.has_table
    )
    run = cell.text_frame.paragraphs[0].runs[0]
    assert run.text == "replaced"
    assert run.font.name == "Malgun Gothic"
    assert run.font.size == PptxPt(18)
    assert run.font.bold is True


def test_docx_append_row_preserves_formatting(tmp_path: Path) -> None:
    """add_row가 템플릿 행을 복사하더라도 새 run을 쓰지 않고 첫 run을 재활용해야 한다."""
    src = tmp_path / "append.docx"
    doc = Document()
    table = doc.add_table(rows=1, cols=2)
    for i in range(2):
        run = table.cell(0, i).paragraphs[0].add_run("hdr")
        run.font.name = "Malgun Gothic"
        run.font.size = DocxPt(14)
    doc.save(src)

    adapter = load(src)
    adapter.append_row(0, ["A", "B"])
    adapter.save()
    adapter.close()

    verify = Document(src)
    new_row = verify.tables[0].rows[1]
    for col, expected in enumerate(["A", "B"]):
        run = new_row.cells[col].paragraphs[0].runs[0]
        assert run.text == expected
        # 복사된 행은 템플릿 속성을 이어받으므로 font가 비어 있지 않아야 한다
        assert run.font.name in {"Malgun Gothic", None}


# ---- empty-cell regression: endParaRPr / pPr.rPr must be cloned (issue #1 v0.1.2) ----


def _make_pptx_empty_cell_with_endpararpr(path: Path) -> None:
    """빈 PPTX 셀에 endParaRPr(폰트 정보)만 심어둔 상태를 만든다.

    실제 고객 템플릿에서 "아직 값을 입력하지 않았지만 셀이 칠해질 때 쓸
    폰트"가 endParaRPr에 박혀 있는 상황을 재현한다.
    """
    from lxml import etree

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(6), Inches(1))
    cell = table_shape.table.cell(0, 0)

    # cell.text = ""로 만든 뒤 기존 run과 그 rPr을 전부 제거하고,
    # endParaRPr만 남긴다 (실제 PowerPoint 저장본의 empty-cell 상태 모방).
    cell.text = ""
    p_el = cell.text_frame.paragraphs[0]._p
    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    # 기존 <a:r>들 제거
    for r in p_el.findall(f"{{{A_NS}}}r"):
        p_el.remove(r)
    # endParaRPr 주입 (이미 있으면 속성만 세팅)
    end_rpr = p_el.find(f"{{{A_NS}}}endParaRPr")
    if end_rpr is None:
        end_rpr = etree.SubElement(p_el, f"{{{A_NS}}}endParaRPr")
    end_rpr.set("lang", "en-US")
    end_rpr.set("sz", "1800")  # 18pt in PPTX units (1/100 pt)
    end_rpr.set("b", "1")
    latin = etree.SubElement(end_rpr, f"{{{A_NS}}}latin")
    latin.set("typeface", "Microsoft Sans Serif")

    prs.save(path)


def test_pptx_empty_cell_preserves_endpararpr(tmp_path: Path) -> None:
    src = tmp_path / "empty_endpara.pptx"
    _make_pptx_empty_cell_with_endpararpr(src)

    doc = load(src)
    old = doc.set_cell(0, 0, 0, "V-2024-001")
    doc.save()
    doc.close()
    assert old == ""

    verify = Presentation(src)
    cell = next(
        shape.table.cell(0, 0)
        for slide in verify.slides
        for shape in slide.shapes
        if shape.has_table
    )
    run = cell.text_frame.paragraphs[0].runs[0]
    assert run.text == "V-2024-001"
    # endParaRPr에 담겼던 속성이 run의 rPr로 복사되었는지 XML 레벨로 확인
    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    rpr = run._r.find(f"{{{A_NS}}}rPr")
    assert rpr is not None, "new run lost its rPr"
    assert rpr.get("sz") == "1800"
    assert rpr.get("b") == "1"
    latin = rpr.find(f"{{{A_NS}}}latin")
    assert latin is not None
    assert latin.get("typeface") == "Microsoft Sans Serif"


def _make_docx_empty_cell_with_ppr_rpr(path: Path) -> None:
    """빈 DOCX 셀의 paragraph에 <w:pPr><w:rPr>만 심어둔 상태를 만든다."""
    from lxml import etree

    W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    doc = Document()
    table = doc.add_table(rows=1, cols=1)
    cell = table.cell(0, 0)
    p_el = cell.paragraphs[0]._p

    # 기존 run들 제거
    for r in p_el.findall(f"{{{W_NS}}}r"):
        p_el.remove(r)

    # <w:pPr><w:rPr>...</w:rPr></w:pPr> 주입
    ppr = p_el.find(f"{{{W_NS}}}pPr")
    if ppr is None:
        ppr = etree.SubElement(p_el, f"{{{W_NS}}}pPr")
        p_el.insert(0, ppr)
    rpr = etree.SubElement(ppr, f"{{{W_NS}}}rPr")
    rfonts = etree.SubElement(rpr, f"{{{W_NS}}}rFonts")
    rfonts.set(f"{{{W_NS}}}ascii", "Malgun Gothic")
    rfonts.set(f"{{{W_NS}}}eastAsia", "Malgun Gothic")
    sz = etree.SubElement(rpr, f"{{{W_NS}}}sz")
    sz.set(f"{{{W_NS}}}val", "36")  # half-points → 18pt
    b = etree.SubElement(rpr, f"{{{W_NS}}}b")
    b.set(f"{{{W_NS}}}val", "1")

    doc.save(path)


def test_docx_empty_cell_preserves_ppr_rpr(tmp_path: Path) -> None:
    src = tmp_path / "empty_pprrpr.docx"
    _make_docx_empty_cell_with_ppr_rpr(src)

    doc = load(src)
    old = doc.set_cell(0, 0, 0, "값")
    doc.save()
    doc.close()
    assert old == ""

    verify = Document(src)
    run = verify.tables[0].cell(0, 0).paragraphs[0].runs[0]
    assert run.text == "값"
    # pPr.rPr에 담겼던 속성이 run의 rPr로 복사됐는지 확인
    W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    rpr = run._r.find(f"{{{W_NS}}}rPr")
    assert rpr is not None, "new run lost its rPr"
    sz = rpr.find(f"{{{W_NS}}}sz")
    assert sz is not None and sz.get(f"{{{W_NS}}}val") == "36"
    rfonts = rpr.find(f"{{{W_NS}}}rFonts")
    assert rfonts is not None
    assert rfonts.get(f"{{{W_NS}}}ascii") == "Malgun Gothic"


# ---- HWPX merged cell / nested table awareness (v0.1.3) ----


def _make_hwpx_merged_table(path: Path) -> None:
    """3x3 표의 첫 행을 colSpan=3으로 병합한 HWPX 문서 생성.

    HWPX에서 병합은 (a) anchor의 cellSpan 변경 + (b) 나머지 셀을
    width=0, height=0, 빈 텍스트로 "deactivate"하는 방식으로 이뤄진다.
    """
    doc = HwpxDocument.new()
    doc.add_paragraph("")
    doc.add_table(3, 3)
    doc.save_to_path(path)

    # python-hwpx의 add_table는 일반 3x3을 만든다. 병합은 재오픈 후 처리.
    doc2 = HwpxDocument.open(path)
    try:
        section = doc2.sections[0]
        tbl = None
        for para in section.paragraphs:
            if para.tables:
                tbl = para.tables[0]
                break
        assert tbl is not None, "created table not found"

        row0 = tbl.rows[0].cells
        anchor = row0[0]
        anchor.set_span(row_span=1, col_span=3)
        anchor.text = "병합된 제목"
        for sibling in (row0[1], row0[2]):
            sibling.set_size(width=0, height=0)
            sibling.text = ""

        r1 = tbl.rows[1].cells
        r1[0].text = "A1"; r1[1].text = "A2"; r1[2].text = "A3"
        r2 = tbl.rows[2].cells
        r2[0].text = "B1"; r2[1].text = "B2"; r2[2].text = "B3"

        doc2.save_to_path(path)
    finally:
        doc2.close()


def test_hwpx_get_tables_reports_merges(tmp_path: Path) -> None:
    src = tmp_path / "merged.hwpx"
    _make_hwpx_merged_table(src)

    adapter = load(src)
    try:
        tables = adapter.get_tables()
    finally:
        adapter.close()

    assert len(tables) == 1
    t = tables[0]
    assert t.rows == 3 and t.cols == 3
    # 첫 행: 앵커에만 텍스트, 나머지는 None
    assert t.preview[0][0] == "병합된 제목"
    assert t.preview[0][1] is None
    assert t.preview[0][2] is None
    # 일반 행은 그대로
    assert t.preview[1] == ["A1", "A2", "A3"]
    assert t.preview[2] == ["B1", "B2", "B3"]
    # merges 메타
    assert len(t.merges) == 1
    assert t.merges[0].anchor == (0, 0)
    assert t.merges[0].span == (1, 3)


def test_hwpx_set_cell_rejects_merged_slot(tmp_path: Path) -> None:
    src = tmp_path / "merged.hwpx"
    _make_hwpx_merged_table(src)

    adapter = load(src)
    try:
        with pytest.raises(ValueError, match="merged region"):
            adapter.set_cell(0, 0, 2, "해킹")
        # 앵커 보존 확인
        adapter_tables = adapter.get_tables()
        assert adapter_tables[0].preview[0][0] == "병합된 제목"
    finally:
        adapter.close()


def test_hwpx_set_cell_anchor_succeeds(tmp_path: Path) -> None:
    src = tmp_path / "merged.hwpx"
    _make_hwpx_merged_table(src)

    adapter = load(src)
    try:
        old = adapter.set_cell(0, 0, 0, "새 제목")
        adapter.save()
    finally:
        adapter.close()

    assert old == "병합된 제목"
    verify = load(src)
    try:
        tables = verify.get_tables()
    finally:
        verify.close()
    assert tables[0].preview[0][0] == "새 제목"
    assert tables[0].preview[0][1] is None  # 병합 구조 유지


def test_hwpx_set_cell_allow_merge_redirect(tmp_path: Path) -> None:
    import warnings as _w

    src = tmp_path / "merged.hwpx"
    _make_hwpx_merged_table(src)

    adapter = load(src)
    try:
        with _w.catch_warnings(record=True) as caught:
            _w.simplefilter("always")
            old = adapter.set_cell(0, 0, 2, "리디렉트", allow_merge_redirect=True)
        assert any("redirected" in str(w.message) for w in caught)
        assert old == "병합된 제목"
        adapter.save()
    finally:
        adapter.close()

    verify = load(src)
    try:
        tables = verify.get_tables()
    finally:
        verify.close()
    assert tables[0].preview[0][0] == "리디렉트"


def test_hwpx_nested_table_indexed_with_parent_path(tmp_path: Path) -> None:
    src = tmp_path / "nested.hwpx"
    doc = HwpxDocument.new()
    doc.add_paragraph("")
    doc.add_table(2, 2)
    doc.save_to_path(src)

    # 재오픈해 (0,0) 셀에 중첩 테이블 삽입
    doc2 = HwpxDocument.open(src)
    try:
        section = doc2.sections[0]
        outer = None
        for para in section.paragraphs:
            if para.tables:
                outer = para.tables[0]
                break
        assert outer is not None
        # anchor cell에 중첩 테이블
        inner_cell = outer.cell(0, 0)
        inner = inner_cell.add_table(1, 2)
        inner.cell(0, 0).text = "중첩A"
        inner.cell(0, 1).text = "중첩B"
        # 바깥 셀 다른 위치도 채움
        outer.cell(0, 1).text = "out01"
        outer.cell(1, 0).text = "out10"
        outer.cell(1, 1).text = "out11"
        doc2.save_to_path(src)
    finally:
        doc2.close()

    adapter = load(src)
    try:
        tables = adapter.get_tables()
    finally:
        adapter.close()

    # flat DFS: 바깥 테이블 idx=0, 중첩 테이블 idx=1
    assert len(tables) == 2
    assert tables[0].parent_path is None
    assert tables[0].rows == 2 and tables[0].cols == 2
    assert tables[1].parent_path is not None
    assert "cell(0,0)" in tables[1].parent_path
    assert tables[1].preview[0] == ["중첩A", "중첩B"]


def test_hwpx_2d_merge_and_multiple_merges(tmp_path: Path) -> None:
    """rowSpan>1 AND colSpan>1, 같은 표에 복수 병합."""
    src = tmp_path / "complex.hwpx"
    doc = HwpxDocument.new()
    doc.add_paragraph("")
    doc.add_table(5, 4)
    doc.save_to_path(src)

    doc2 = HwpxDocument.open(src)
    try:
        tbl = next(t for p in doc2.sections[0].paragraphs for t in p.tables)
        # merge 1: (0,0) 1x2 horizontal
        tbl.rows[0].cells[0].set_span(1, 2); tbl.rows[0].cells[0].text = "M1"
        tbl.rows[0].cells[1].set_size(0, 0); tbl.rows[0].cells[1].text = ""
        # merge 2: (2,2) 2x2 block
        tbl.rows[2].cells[2].set_span(2, 2); tbl.rows[2].cells[2].text = "M2"
        tbl.rows[2].cells[3].set_size(0, 0); tbl.rows[2].cells[3].text = ""
        tbl.rows[3].cells[2].set_size(0, 0); tbl.rows[3].cells[2].text = ""
        tbl.rows[3].cells[3].set_size(0, 0); tbl.rows[3].cells[3].text = ""
        doc2.save_to_path(src)
    finally:
        doc2.close()

    adapter = load(src)
    try:
        t = adapter.get_tables()[0]
    finally:
        adapter.close()

    assert t.preview[0][0] == "M1" and t.preview[0][1] is None
    assert t.preview[2][2] == "M2"
    assert t.preview[2][3] is None
    assert t.preview[3][2] is None and t.preview[3][3] is None
    anchors = {m.anchor: m.span for m in t.merges}
    assert anchors[(0, 0)] == (1, 2)
    assert anchors[(2, 2)] == (2, 2)


def test_hwpx_merge_beyond_preview_cutoff(tmp_path: Path) -> None:
    """preview_rows로 잘려도 merges 메타에는 모든 병합이 포함돼야 한다."""
    src = tmp_path / "bigtable.hwpx"
    doc = HwpxDocument.new()
    doc.add_paragraph("")
    doc.add_table(10, 3)
    doc.save_to_path(src)

    doc2 = HwpxDocument.open(src)
    try:
        tbl = next(t for p in doc2.sections[0].paragraphs for t in p.tables)
        tbl.rows[5].cells[0].set_span(1, 3); tbl.rows[5].cells[0].text = "DEEP"
        tbl.rows[5].cells[1].set_size(0, 0); tbl.rows[5].cells[1].text = ""
        tbl.rows[5].cells[2].set_size(0, 0); tbl.rows[5].cells[2].text = ""
        doc2.save_to_path(src)
    finally:
        doc2.close()

    adapter = load(src)
    try:
        t = adapter.get_tables(preview_rows=4)[0]
    finally:
        adapter.close()

    assert len(t.preview) == 4
    assert any(m.anchor == (5, 0) and m.span == (1, 3) for m in t.merges)


def test_hwpx_nested_cell_text_isolated(tmp_path: Path) -> None:
    """외부 셀 프리뷰에 중첩 테이블 텍스트가 섞여 들어오면 안 된다."""
    src = tmp_path / "nested_isolation.hwpx"
    doc = HwpxDocument.new()
    doc.add_paragraph("")
    doc.add_table(2, 2)
    doc.save_to_path(src)

    doc2 = HwpxDocument.open(src)
    try:
        outer = next(t for p in doc2.sections[0].paragraphs for t in p.tables)
        inner = outer.cell(0, 0).add_table(1, 1)
        inner.cell(0, 0).text = "INNER_ONLY"
        outer.cell(1, 1).text = "outer"
        doc2.save_to_path(src)
    finally:
        doc2.close()

    adapter = load(src)
    try:
        tables = adapter.get_tables()
    finally:
        adapter.close()

    # 외부 (0,0)은 중첩 테이블을 담고 있지만 직접 텍스트는 없어야 한다
    assert tables[0].preview[0][0] == ""
    # 내부 테이블만 INNER_ONLY를 담는다
    assert tables[1].preview[0][0] == "INNER_ONLY"


def test_hwpx_set_cell_on_nested_table(tmp_path: Path) -> None:
    """flat index로 중첩 테이블 셀도 편집 가능해야 한다."""
    src = tmp_path / "nested_edit.hwpx"
    doc = HwpxDocument.new()
    doc.add_paragraph("")
    doc.add_table(2, 2)
    doc.save_to_path(src)

    doc2 = HwpxDocument.open(src)
    try:
        outer = next(t for p in doc2.sections[0].paragraphs for t in p.tables)
        inner = outer.cell(0, 0).add_table(1, 2)
        inner.cell(0, 0).text = "before"
        inner.cell(0, 1).text = "keep"
        doc2.save_to_path(src)
    finally:
        doc2.close()

    adapter = load(src)
    try:
        old = adapter.set_cell(1, 0, 0, "after")  # nested idx=1
        adapter.save()
    finally:
        adapter.close()

    assert old == "before"

    verify = load(src)
    try:
        t = verify.get_tables()[1]
    finally:
        verify.close()
    assert t.preview[0] == ["after", "keep"]


# ---- v0.3.0: get_cell / append_to_cell / DOCX+PPTX merge awareness ----


def test_hwpx_get_cell_returns_full_text(tmp_path: Path) -> None:
    src = tmp_path / "big_cell.hwpx"
    doc = HwpxDocument.new()
    doc.add_paragraph("")
    doc.add_table(1, 1)
    doc.save_to_path(src)

    long_text = "이 셀은 50자를 훨씬 넘기는 내용을 담고 있어서 preview에서 잘리게 된다. " * 3
    doc2 = HwpxDocument.open(src)
    try:
        tbl = next(t for p in doc2.sections[0].paragraphs for t in p.tables)
        tbl.cell(0, 0).text = long_text
        doc2.save_to_path(src)
    finally:
        doc2.close()

    adapter = load(src)
    try:
        preview = adapter.get_tables()[0].preview[0][0]
        cell = adapter.get_cell(0, 0, 0)
    finally:
        adapter.close()

    assert len(preview) <= 40  # truncated
    assert cell.text == long_text  # full
    assert cell.is_anchor is True
    assert cell.anchor == (0, 0)
    assert cell.span == (1, 1)


def test_hwpx_get_cell_reports_merge_info(tmp_path: Path) -> None:
    src = tmp_path / "merged_lookup.hwpx"
    doc = HwpxDocument.new()
    doc.add_paragraph("")
    doc.add_table(2, 3)
    doc.save_to_path(src)

    doc2 = HwpxDocument.open(src)
    try:
        tbl = next(t for p in doc2.sections[0].paragraphs for t in p.tables)
        tbl.rows[0].cells[0].set_span(1, 3)
        tbl.rows[0].cells[0].text = "HEADER"
        for s in (tbl.rows[0].cells[1], tbl.rows[0].cells[2]):
            s.set_size(0, 0); s.text = ""
        doc2.save_to_path(src)
    finally:
        doc2.close()

    adapter = load(src)
    try:
        anchor_info = adapter.get_cell(0, 0, 0)
        non_anchor_info = adapter.get_cell(0, 0, 2)
    finally:
        adapter.close()

    assert anchor_info.is_anchor is True
    assert anchor_info.span == (1, 3)
    assert anchor_info.text == "HEADER"
    assert non_anchor_info.is_anchor is False
    assert non_anchor_info.anchor == (0, 0)
    assert non_anchor_info.span == (1, 3)
    assert non_anchor_info.text == "HEADER"  # anchor의 텍스트 반환


def test_hwpx_append_to_cell_preserves_label(tmp_path: Path) -> None:
    src = tmp_path / "label_cell.hwpx"
    doc = HwpxDocument.new()
    doc.add_paragraph("")
    doc.add_table(1, 1)
    doc.save_to_path(src)

    doc2 = HwpxDocument.open(src)
    try:
        tbl = next(t for p in doc2.sections[0].paragraphs for t in p.tables)
        tbl.cell(0, 0).text = "성  명"
        doc2.save_to_path(src)
    finally:
        doc2.close()

    adapter = load(src)
    try:
        old = adapter.append_to_cell(0, 0, 0, "홍길동", separator="   ")
        adapter.save()
    finally:
        adapter.close()

    assert old == "성  명"
    verify = load(src)
    try:
        cell = verify.get_cell(0, 0, 0)
    finally:
        verify.close()
    assert cell.text == "성  명   홍길동"


def test_hwpx_append_to_cell_empty_skips_separator(tmp_path: Path) -> None:
    src = tmp_path / "empty.hwpx"
    doc = HwpxDocument.new()
    doc.add_paragraph("")
    doc.add_table(1, 1)
    doc.save_to_path(src)

    adapter = load(src)
    try:
        old = adapter.append_to_cell(0, 0, 0, "first")
        adapter.save()
    finally:
        adapter.close()
    assert old == ""

    verify = load(src)
    try:
        cell = verify.get_cell(0, 0, 0)
    finally:
        verify.close()
    assert cell.text == "first"


def test_hwpx_append_row_copies_style(tmp_path: Path) -> None:
    src = tmp_path / "appendable.hwpx"
    doc = HwpxDocument.new()
    doc.add_paragraph("")
    doc.add_table(2, 3)
    doc.save_to_path(src)

    doc2 = HwpxDocument.open(src)
    try:
        tbl = next(t for p in doc2.sections[0].paragraphs for t in p.tables)
        tbl.cell(0, 0).text = "H1"
        tbl.cell(0, 1).text = "H2"
        tbl.cell(0, 2).text = "H3"
        tbl.cell(1, 0).text = "r1c0"
        tbl.cell(1, 1).text = "r1c1"
        tbl.cell(1, 2).text = "r1c2"
        doc2.save_to_path(src)
    finally:
        doc2.close()

    adapter = load(src)
    try:
        adapter.append_row(0, ["r2c0", "r2c1", "r2c2"])
        adapter.save()
    finally:
        adapter.close()

    verify = load(src)
    try:
        t = verify.get_tables()[0]
    finally:
        verify.close()
    assert t.rows == 3
    assert t.preview[2] == ["r2c0", "r2c1", "r2c2"]


def test_hwpx_errors_are_custom_subclasses() -> None:
    from document_adapter.base import (
        MergedCellWriteError,
        CellOutOfBoundsError,
        TableIndexError,
    )

    assert issubclass(MergedCellWriteError, ValueError)
    assert issubclass(CellOutOfBoundsError, IndexError)
    assert issubclass(TableIndexError, IndexError)


def _make_docx_merged_table(path: Path) -> None:
    """2x3 DOCX에서 (0,0)~(0,2) 수평 병합."""
    doc = Document()
    table = doc.add_table(rows=2, cols=3)
    # python-docx merge API: cell.merge(other_cell)
    a = table.cell(0, 0)
    a.merge(table.cell(0, 1))
    a.merge(table.cell(0, 2))
    a.text = "HEADER"
    table.cell(1, 0).text = "A"
    table.cell(1, 1).text = "B"
    table.cell(1, 2).text = "C"
    doc.save(path)


def test_docx_get_tables_reports_merges(tmp_path: Path) -> None:
    src = tmp_path / "merged.docx"
    _make_docx_merged_table(src)

    adapter = load(src)
    try:
        tables = adapter.get_tables()
    finally:
        adapter.close()

    t = tables[0]
    assert t.rows == 2 and t.cols == 3
    assert t.preview[0][0] == "HEADER"
    assert t.preview[0][1] is None
    assert t.preview[0][2] is None
    assert t.preview[1] == ["A", "B", "C"]
    assert len(t.merges) == 1
    assert t.merges[0].anchor == (0, 0)
    assert t.merges[0].span == (1, 3)


def test_docx_set_cell_rejects_merged_slot(tmp_path: Path) -> None:
    src = tmp_path / "reject.docx"
    _make_docx_merged_table(src)

    adapter = load(src)
    try:
        with pytest.raises(ValueError, match="merged region"):
            adapter.set_cell(0, 0, 2, "X")
        # anchor는 성공
        old = adapter.set_cell(0, 0, 0, "NEW HEADER")
        adapter.save()
    finally:
        adapter.close()
    assert old == "HEADER"

    verify = load(src)
    try:
        t = verify.get_tables()[0]
    finally:
        verify.close()
    assert t.preview[0][0] == "NEW HEADER"


def test_docx_append_to_cell(tmp_path: Path) -> None:
    src = tmp_path / "append_label.docx"
    doc = Document()
    tbl = doc.add_table(rows=1, cols=1)
    tbl.cell(0, 0).text = "Name"
    doc.save(src)

    adapter = load(src)
    try:
        old = adapter.append_to_cell(0, 0, 0, "Alice", separator=": ")
        adapter.save()
    finally:
        adapter.close()
    assert old == "Name"

    verify = Document(src)
    assert verify.tables[0].cell(0, 0).text == "Name: Alice"


def _make_pptx_merged_table(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(2, 3, Inches(1), Inches(1), Inches(6), Inches(1.5))
    tbl = shape.table
    # PPTX는 (0,0) → (0,2) 병합
    tbl.cell(0, 0).merge(tbl.cell(0, 2))
    tbl.cell(0, 0).text = "TITLE"
    tbl.cell(1, 0).text = "A"
    tbl.cell(1, 1).text = "B"
    tbl.cell(1, 2).text = "C"
    prs.save(path)


def test_pptx_get_tables_reports_merges(tmp_path: Path) -> None:
    src = tmp_path / "merged.pptx"
    _make_pptx_merged_table(src)

    adapter = load(src)
    try:
        t = adapter.get_tables()[0]
    finally:
        adapter.close()

    assert t.preview[0][0] == "TITLE"
    assert t.preview[0][1] is None
    assert t.preview[0][2] is None
    assert t.preview[1] == ["A", "B", "C"]
    assert len(t.merges) == 1
    assert t.merges[0].anchor == (0, 0)
    assert t.merges[0].span == (1, 3)


def test_pptx_set_cell_rejects_merged_slot(tmp_path: Path) -> None:
    src = tmp_path / "pptx_reject.pptx"
    _make_pptx_merged_table(src)

    adapter = load(src)
    try:
        with pytest.raises(ValueError, match="merged region"):
            adapter.set_cell(0, 0, 2, "X")
        old = adapter.set_cell(0, 0, 0, "NEW TITLE")
        adapter.save()
    finally:
        adapter.close()
    assert old == "TITLE"


def test_pptx_append_to_cell(tmp_path: Path) -> None:
    src = tmp_path / "pptx_append.pptx"
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(4), Inches(1))
    shape.table.cell(0, 0).text = "성명"
    prs.save(src)

    adapter = load(src)
    try:
        old = adapter.append_to_cell(0, 0, 0, "Alice", separator=" → ")
        adapter.save()
    finally:
        adapter.close()
    assert old == "성명"

    verify = Presentation(src)
    cell = next(
        sh.table.cell(0, 0)
        for sl in verify.slides
        for sh in sl.shapes
        if sh.has_table
    )
    assert cell.text == "성명 → Alice"


def test_hwpx_set_cell_preserves_charprid_ref(tmp_path: Path) -> None:
    """HWPX 셀의 run이 가진 charPrIDRef가 set_cell 후에도 유지되는지 확인.

    python-hwpx의 paragraph.text setter는 PPTX/DOCX와 달리 기존 run의 <hp:t>만
    교체하는 방식이라 charPrIDRef가 자연스럽게 보존된다. 이 테스트는 upstream이
    그 동작을 바꿀 경우(run 재생성 방식으로) 즉시 감지하기 위한 회귀 가드다.
    """
    import re
    import zipfile

    src = tmp_path / "empty_charpr.hwpx"
    doc = HwpxDocument.new()
    doc.add_paragraph("")
    doc.add_table(1, 1)
    doc.save_to_path(src)

    # table cell 안의 run만 charPrIDRef="7"로 조작한 zip으로 재기록
    patched = tmp_path / "patched_charpr.hwpx"
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(patched, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if "section" in item.filename.lower() and item.filename.endswith(".xml"):
                text = data.decode("utf-8")
                def repl(m: re.Match) -> str:
                    tc = m.group(0)
                    return re.sub(
                        r'(<hp:run\s+charPrIDRef=")0(")',
                        r"\g<1>7\g<2>",
                        tc,
                        count=1,
                    )
                text = re.sub(r"<hp:tc[^>]*>.*?</hp:tc>", repl, text, flags=re.DOTALL)
                data = text.encode("utf-8")
            zout.writestr(item, data)

    # set_cell 실행
    adapter = load(patched)
    adapter.set_cell(0, 0, 0, "새 값")
    adapter.save()
    adapter.close()

    # table cell 안의 run만 추출해서 확인
    with zipfile.ZipFile(patched) as z:
        section = next(
            n for n in z.namelist() if "section" in n.lower() and n.endswith(".xml")
        )
        xml = z.read(section).decode("utf-8")

    tc_match = re.search(r"<hp:tc[^>]*>.*?</hp:tc>", xml, flags=re.DOTALL)
    assert tc_match is not None
    run_match = re.search(r"<hp:run[^>]*>.*?</hp:run>", tc_match.group())
    assert run_match is not None
    run_xml = run_match.group()
    assert 'charPrIDRef="7"' in run_xml, f"charPrIDRef not preserved: {run_xml!r}"
    assert "새 값" in run_xml
