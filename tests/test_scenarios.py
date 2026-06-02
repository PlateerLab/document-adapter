"""시나리오 기반 회귀/보안 테스트.

test_smoke.py 가 "정상 동작"을 검증한다면, 이 파일은 발굴한 엣지케이스·보안
시나리오를 실패 테스트로 고정하고 개선을 추적한다.

    pytest tests/test_scenarios.py -v
"""
from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from lxml import etree
from hwpx.document import HwpxDocument

from document_adapter import load
from document_adapter.hwpx_core.package import HwpxPackage


# ---------------------------------------------------------------------------
# 공용 헬퍼
# ---------------------------------------------------------------------------

def _make_hwpx(path: Path, paragraphs: list[str]) -> None:
    doc = HwpxDocument.new()
    for p in paragraphs:
        doc.add_paragraph(p)
    doc.save_to_path(path)


def _poison_section0(src: Path, dst: Path, *, doctype: bytes, marker: bytes) -> None:
    """src.hwpx 의 Contents/section0.xml 에 DOCTYPE 엔티티를 주입한다.

    - doctype: <!DOCTYPE ...> 선언 (엔티티 정의 포함)
    - marker: 본문 어딘가의 바이트를 엔티티 참조로 치환할 (orig, ref) 쌍
    나머지 파트는 그대로 복사해 유효한 .hwpx 컨테이너를 유지한다.
    """
    with zipfile.ZipFile(src, "r") as zin:
        names = zin.namelist()
        raw = {n: zin.read(n) for n in names}

    sec = raw["Contents/section0.xml"]
    # XML 선언 뒤에 DOCTYPE 삽입
    assert sec.lstrip().startswith(b"<?xml"), "예상한 XML 선언이 없음"
    decl_end = sec.index(b"?>") + 2
    sec = sec[:decl_end] + b"\n" + doctype + sec[decl_end:]
    # 본문에 엔티티 참조를 끼워 넣을 자리: 첫 <hp:t> ... </hp:t> 사이에 marker 주입
    orig, ref = marker
    sec = sec.replace(orig, orig + ref, 1)
    raw["Contents/section0.xml"] = sec

    with zipfile.ZipFile(dst, "w") as zout:
        for n in names:
            zout.writestr(n, raw[n])


# ---------------------------------------------------------------------------
# 시나리오 A — XXE / 엔티티 확장 방어 (보안)
# ---------------------------------------------------------------------------

def test_hwpx_internal_entity_not_resolved(tmp_path: Path) -> None:
    """신뢰할 수 없는 .hwpx 의 내부 엔티티가 확장되면 안 된다.

    하드닝되지 않은 파서(resolve_entities=True 기본값)는 &pwn; 을 'PWNED'로
    확장한다. 보안상 untrusted 입력의 엔티티는 절대 확장하면 안 된다.
    """
    good = tmp_path / "good.hwpx"
    _make_hwpx(good, ["HELLO"])
    bad = tmp_path / "bad.hwpx"
    _poison_section0(
        good, bad,
        doctype=b'<!DOCTYPE hs:sec [ <!ENTITY pwn "PWNED"> ]>',
        marker=(b"HELLO", b"&pwn;"),
    )

    pkg = HwpxPackage.open(bad)
    text = pkg.export_text()
    pkg.close()
    assert "PWNED" not in text, "내부 엔티티가 확장됨 — 파서 하드닝 필요"


def test_hwpx_billion_laughs_not_expanded(tmp_path: Path) -> None:
    """billion-laughs 류 엔티티 확장 DoS 방어."""
    good = tmp_path / "good.hwpx"
    _make_hwpx(good, ["BOOM"])
    bad = tmp_path / "bad.hwpx"
    doctype = (
        b'<!DOCTYPE hs:sec [\n'
        b' <!ENTITY a "aaaaaaaaaa">\n'
        b' <!ENTITY b "&a;&a;&a;&a;&a;&a;&a;&a;&a;&a;">\n'
        b' <!ENTITY c "&b;&b;&b;&b;&b;&b;&b;&b;&b;&b;">\n'
        b' <!ENTITY d "&c;&c;&c;&c;&c;&c;&c;&c;&c;&c;">\n'
        b']>'
    )
    _poison_section0(good, bad, doctype=doctype, marker=(b"BOOM", b"&d;"))

    pkg = HwpxPackage.open(bad)
    text = pkg.export_text()
    pkg.close()
    # 확장됐다면 'a'가 수천 자. 하드닝되면 0자.
    assert text.count("a") < 100, "엔티티 확장 DoS 방어 필요"


# ---------------------------------------------------------------------------
# 시나리오 B — ABC 계약 일치 (allow_merge_redirect)
# ---------------------------------------------------------------------------

def test_abc_signature_supports_fill_form() -> None:
    """DocumentAdapter ABC 의 추상 시그니처만 보고 구현한 어댑터로도
    fill_form 이 동작해야 한다.

    base.fill_form → _fill_one_cell 은 set_cell/append_to_cell 을
    allow_merge_redirect= 키워드로 호출한다. 추상 시그니처에 이 인자가
    선언돼 있지 않으면, ABC 계약대로 구현한 어댑터에서 TypeError 가 난다.
    """
    from document_adapter.base import (
        DocumentAdapter,
        TableSchema,
        CellContent,
    )

    class FakeAdapter(DocumentAdapter):
        """ABC 추상 시그니처 '그대로' 구현한 최소 어댑터 (단일 라벨-값 표)."""
        format = "fake"

        def _open(self) -> None:
            # 2x2: (0,0)="성명" 라벨, (0,1)=빈 값셀
            self._cells = {(0, 0): "성명", (0, 1): "", (1, 0): "생년", (1, 1): ""}

        def save(self, path=None):
            return self.path

        def get_placeholders(self):
            return []

        def get_tables(self, min_rows=1, min_cols=1, preview_rows=4, max_cell_len=40):
            preview = [["성명", ""], ["생년", ""]]
            return [TableSchema(index=0, rows=2, cols=2, preview=preview)]

        def get_cell(self, table_index, row, col):
            text = self._cells.get((row, col), "")
            return CellContent(
                row=row, col=col, text=text, paragraphs=[text],
                is_anchor=True, anchor=(row, col), span=(1, 1),
            )

        def render_template(self, context):
            pass

        # 핵심: ABC 추상 시그니처를 '그대로' 따른다 (allow_merge_redirect 없음)
        def set_cell(self, table_index, row, col, value):
            old = self._cells.get((row, col), "")
            self._cells[(row, col)] = value
            return old

        def append_to_cell(self, table_index, row, col, value, separator="  "):
            old = self._cells.get((row, col), "")
            self._cells[(row, col)] = (old + separator + value) if old else value
            return old

        def append_row(self, table_index, values):
            pass

    fa = FakeAdapter(Path("/dev/null"))
    # 추상 시그니처대로 구현했는데 fill_form 이 깨지면 ABC 계약이 거짓.
    result = fa.fill_form({"성명": "홍길동"})
    assert result["filled"], f"fill_form 이 라벨을 채우지 못함: {result}"
    assert fa._cells[(0, 1)] == "홍길동"


# ---------------------------------------------------------------------------
# 시나리오 C — 발굴한 보장(guarantee)을 회귀 테스트로 고정
# ---------------------------------------------------------------------------

def _make_form_hwpx(path: Path, rows: list[tuple[str, str]]) -> None:
    """라벨-값 2열 표를 가진 HWPX 폼 생성."""
    doc = HwpxDocument.new()
    doc.add_paragraph("")
    doc.add_table(len(rows), 2)
    doc.save_to_path(path)
    d = HwpxDocument.open(path)
    try:
        sec = d.sections[0]
        tbl = next(p.tables[0] for p in sec.paragraphs if p.tables)
        for i, (label, value) in enumerate(rows):
            tbl.rows[i].cells[0].text = label
            tbl.rows[i].cells[1].text = value
        d.save_to_path(path)
    finally:
        d.close()


def test_set_cell_escapes_xml_special_chars(tmp_path: Path) -> None:
    """셀 값에 XML 특수문자/태그 문자열이 들어와도 이스케이프되어
    문서 구조를 깨뜨리지 않고 그대로 보존돼야 한다 (injection 방어)."""
    src = tmp_path / "f.hwpx"
    _make_form_hwpx(src, [("제목", ""), ("비고", "")])
    payload = '<hp:t>x</hp:t> & "q" < > 이모지😀'

    ad = load(src)
    ad.set_cell(0, 0, 1, payload)
    out = tmp_path / "out.hwpx"
    ad.save(out)
    ad.close()

    # well-formed XML 유지
    with zipfile.ZipFile(out) as z:
        etree.fromstring(z.read("Contents/section0.xml"))  # noqa: S320 (신뢰 입력)
    # 값 그대로 보존 + 인접 셀 무손상
    ad2 = load(out)
    assert ad2.get_cell(0, 0, 1).text == payload
    assert ad2.get_cell(0, 1, 0).text == "비고"
    ad2.close()


def test_save_preserves_mimetype_first_and_stored(tmp_path: Path) -> None:
    """save() 후에도 mimetype 이 첫 엔트리 + STORED(비압축)로 유지돼야 한다
    (OPC/HWPX 컨테이너 유효성). 안 건드린 파트는 byte-identical."""
    src = tmp_path / "src.hwpx"
    _make_form_hwpx(src, [("a", "b")])

    pkg = HwpxPackage.open(src)
    pkg.get_xml_root("Contents/section0.xml")
    pkg.mark_dirty("Contents/section0.xml")
    out = tmp_path / "out.hwpx"
    pkg.save(out)
    pkg.close()

    with zipfile.ZipFile(src) as z1, zipfile.ZipFile(out) as z2:
        n2 = z2.namelist()
        assert n2[0] == "mimetype", "mimetype 이 첫 엔트리가 아님"
        assert z2.getinfo("mimetype").compress_type == zipfile.ZIP_STORED
        assert z1.read("mimetype") == z2.read("mimetype")
        # header.xml 은 편집 안 했으므로 byte-identical
        assert z1.read("Contents/header.xml") == z2.read("Contents/header.xml")


@pytest.mark.parametrize("data,expect_filled,expect_notfound,expect_ambig", [
    ({}, 0, 0, 0),                          # 빈 입력
    ({"   ": "x", "***": "y"}, 0, 2, 0),     # 정규화하면 빈 라벨 → not_found
])
def test_fill_form_degenerate_inputs(tmp_path: Path, data, expect_filled,
                                     expect_notfound, expect_ambig) -> None:
    """빈/특수문자뿐인 라벨 등 비정상 입력에서도 crash 없이 분류만."""
    src = tmp_path / "f.hwpx"
    _make_form_hwpx(src, [("성명", ""), ("생년월일", "")])
    ad = load(src)
    r = ad.fill_form(data)
    ad.close()
    assert len(r["filled"]) == expect_filled
    assert len(r["not_found"]) == expect_notfound
    assert len(r["ambiguous"]) == expect_ambig


def test_fill_form_duplicate_label_is_ambiguous(tmp_path: Path) -> None:
    """같은 라벨이 여러 곳이면 채우지 않고 ambiguous 로 분류 (오기입 방지)."""
    src = tmp_path / "dup.hwpx"
    _make_form_hwpx(src, [("금액", ""), ("금액", "")])
    ad = load(src)
    r = ad.fill_form({"금액": "100"})
    ad.close()
    assert len(r["ambiguous"]) == 1
    assert len(r["filled"]) == 0


def _make_hwpx_with_form_controls(path: Path) -> None:
    """python-hwpx 스켈레톤 section0 에 checkBtn/edit 폼 컨트롤을 주입."""
    base = path.with_suffix(".base.hwpx")
    _make_form_hwpx(base, [("동의", "")])
    HP = "http://www.hancom.co.kr/hwpml/2011/paragraph"
    with zipfile.ZipFile(base, "r") as zin:
        names = zin.namelist()
        raw = {n: zin.read(n) for n in names}
    root = etree.fromstring(raw["Contents/section0.xml"])
    p = root.find(f"{{{HP}}}p")
    run = etree.SubElement(p, f"{{{HP}}}run")
    cb = etree.SubElement(run, f"{{{HP}}}checkBtn")
    cb.set("name", "CB1")
    cb.set("caption", "동의")
    cb.set("value", "UNCHECKED")
    ed = etree.SubElement(run, f"{{{HP}}}edit")
    ed.set("name", "ED1")
    etree.SubElement(ed, f"{{{HP}}}text")
    cbx = etree.SubElement(run, f"{{{HP}}}comboBox")
    cbx.set("name", "CBX1")
    for opt in ("서울", "부산"):
        li = etree.SubElement(cbx, f"{{{HP}}}listItem")
        li.set("displayText", opt)
        li.set("value", opt)
    raw["Contents/section0.xml"] = etree.tostring(root, xml_declaration=True,
                                                  encoding="UTF-8")
    with zipfile.ZipFile(path, "w") as zout:
        for n in names:
            zout.writestr(n, raw[n])
    base.unlink()


def test_form_controls_get_and_set(tmp_path: Path) -> None:
    """HWPX 폼 컨트롤(체크박스/에디트) 읽기·쓰기·영속 (#1 갭)."""
    src = tmp_path / "fc.hwpx"
    _make_hwpx_with_form_controls(src)
    ad = load(src)
    ctrls = {c["name"]: c for c in ad.get_form_controls()}
    assert "CB1" in ctrls and ctrls["CB1"]["kind"] == "checkBtn"
    assert "ED1" in ctrls and ctrls["ED1"]["kind"] == "edit"
    assert ctrls["CB1"]["checked"] is False

    assert ctrls["CBX1"]["kind"] == "comboBox"
    assert ctrls["CBX1"]["items"] == ["서울", "부산"]   # 드롭다운 옵션 노출

    ad.set_form_control("CB1", "Y")
    ad.set_form_control("ED1", "홍길동")
    ad.set_form_control("CBX1", "부산")
    ad.save(src)
    ad.close()

    ad2 = load(src)
    after = {c["name"]: c for c in ad2.get_form_controls()}
    ad2.close()
    assert after["CB1"]["checked"] is True
    assert after["ED1"]["value"] == "홍길동"
    assert after["CBX1"]["value"] == "부산"          # 콤보 현재값 영속


def test_form_controls_unsupported_format(tmp_path: Path) -> None:
    """DOCX/PPTX 는 폼 컨트롤 기본 빈 리스트 + set 시 NotImplementedForFormat."""
    from document_adapter.base import NotImplementedForFormat

    from document_adapter.eval.fixtures import docx_blank_form

    src = tmp_path / "f.docx"
    docx_blank_form(src)
    ad = load(src)
    assert ad.get_form_controls() == []
    with pytest.raises(NotImplementedForFormat):
        ad.set_form_control("x", "y")
    ad.close()


def _make_docx_irregular_merge(path: Path) -> None:
    """위 행에 정렬된 tc 가 없는 vMerge=continue 를 가진 docx 표를 만든다.

    이 패턴은 python-docx 의 ``row.cells`` 를 ValueError 로 깨뜨린다
    (tc_at_grid_offset). document-adapter 가 OOXML 레이어에서 그리드를
    직접 계산해 견뎌야 한다.
    """
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document()
    t = doc.add_table(rows=2, cols=4)
    tr0 = t.rows[0]._tr
    tcs0 = tr0.findall(qn("w:tc"))
    p0 = tcs0[0].get_or_add_tcPr()
    p0.append(p0.makeelement(qn("w:gridSpan"), {qn("w:val"): "4"}))
    p0.append(p0.makeelement(qn("w:vMerge"), {qn("w:val"): "restart"}))
    for tc in tcs0[1:]:
        tr0.remove(tc)
    tr1 = t.rows[1]._tr
    tcs1 = tr1.findall(qn("w:tc"))
    pa = tcs1[0].get_or_add_tcPr()
    pa.append(pa.makeelement(qn("w:gridSpan"), {qn("w:val"): "2"}))
    pb = tcs1[1].get_or_add_tcPr()
    pb.append(pb.makeelement(qn("w:gridSpan"), {qn("w:val"): "2"}))
    pb.append(pb.makeelement(qn("w:vMerge"), {qn("w:val"): "continue"}))
    for tc in tcs1[2:]:
        tr1.remove(tc)
    doc.save(str(path))


def test_docx_irregular_horizontal_merge_loads(tmp_path: Path) -> None:
    """가로(gridSpan)+세로(vMerge) 병합이 섞인 docx 표가 깨지지 않고 로딩돼야 한다.

    회귀: python-docx 의 row.cells 가 ValueError 로 깨지던 실제 docx 폼
    (docxtpl horizontal_merge_tpl)에서 발견.
    """
    src = tmp_path / "m.docx"
    _make_docx_irregular_merge(src)
    ad = load(src)
    tbls = ad.get_tables()        # 수정 전에는 여기서 ValueError 로 크래시
    ad_cols = tbls[0].cols
    cell = ad.get_cell(0, 0, 0)   # 셀 접근도 크래시 없이 가능해야
    placeholders = ad.get_placeholders()   # get_placeholders 도 row.cells 미사용
    schema_fmt = ad.get_schema().format    # inspect_document 경로 전체 무크래시
    ad.close()
    assert tbls and ad_cols == 4
    assert cell.span[1] == 4      # gridSpan=4 → colspan 4 (가로병합 인식)
    assert isinstance(placeholders, list)
    assert schema_fmt == "docx"


def _make_template(path: Path, fmt: str) -> None:
    """[{{name}}|{{missing}}] 한 줄 템플릿을 fmt 포맷으로 생성."""
    txt = "[{{name}}|{{missing}}]"
    if fmt == "docx":
        from docx import Document
        d = Document()
        d.add_paragraph(txt)
        d.save(str(path))
    elif fmt == "pptx":
        from pptx import Presentation
        from pptx.util import Inches
        pr = Presentation()
        pr.slide_width = Inches(10)
        pr.slide_height = Inches(7.5)
        s = pr.slides.add_slide(pr.slide_layouts[6])
        s.shapes.add_textbox(Inches(1), Inches(1), Inches(8),
                             Inches(1)).text_frame.text = txt
        pr.save(str(path))
    else:
        from hwpx.document import HwpxDocument
        d2 = HwpxDocument.new()
        d2.add_paragraph(txt)
        d2.save_to_path(path)


def _rendered_text(path: Path) -> str:
    with zipfile.ZipFile(path) as z:
        out = []
        for n in z.namelist():
            if not n.endswith(".xml"):
                continue
            try:
                root = etree.fromstring(z.read(n))
            except etree.XMLSyntaxError:
                continue
            for t in root.iter():
                if t.text and "[" in t.text:
                    out.append(t.text)
        return " ".join(out)


@pytest.mark.parametrize("fmt", ["docx", "pptx", "hwpx"])
def test_render_template_missing_key_blanks_consistently(tmp_path: Path,
                                                         fmt: str) -> None:
    """3포맷 모두 누락 키를 빈칸으로(기본) — {{missing}} 리터럴이 출력에 남지 않음.

    회귀: 이전엔 pptx/hwpx 가 {{missing}} 을 그대로 노출(docx 는 빈칸) — 불일치.
    """
    src = tmp_path / f"t.{fmt}"
    _make_template(src, fmt)
    ad = load(src)
    report = ad.render_template({"name": "홍길동"})   # 기본 on_missing="blank"
    ad.save(src)
    ad.close()
    assert report["used"] == ["name"]
    assert report["missing"] == ["missing"]
    text = _rendered_text(src)
    assert "홍길동" in text
    assert "{{" not in text   # 미완성 플레이스홀더 노출 없음


@pytest.mark.parametrize("fmt", ["docx", "pptx", "hwpx"])
def test_render_template_on_missing_modes(tmp_path: Path, fmt: str) -> None:
    """on_missing leave 는 {{key}} 유지, error 는 ValueError."""
    src = tmp_path / f"t.{fmt}"
    _make_template(src, fmt)
    ad = load(src)
    ad.render_template({"name": "홍길동"}, on_missing="leave")
    ad.save(src)
    leave_text = _rendered_text(src)
    ad.close()
    assert "missing" in leave_text and "}}" in leave_text

    src2 = tmp_path / f"e.{fmt}"
    _make_template(src2, fmt)
    ad2 = load(src2)
    with pytest.raises(ValueError):
        ad2.render_template({"name": "x"}, on_missing="error")
    ad2.close()


def test_overflow_risk_helper() -> None:
    """폭 추정 기반 오버플로 위험 판정: 좁은 칸=위험, 넓은 칸=안전, 미상=보류."""
    from document_adapter.base import _overflow_risk

    assert _overflow_risk("2026-06-01", 0.4) is True    # 0.4cm 스페이서 → 깨짐
    assert _overflow_risk("2026-06-01", 6.4) is False    # 6.4cm 값칸 → 정상
    assert _overflow_risk("서울시 강남구 테헤란로 1", 2.0) is True
    assert _overflow_risk("x", None) is False            # 폭 미상 → 판단 보류
    assert _overflow_risk("", 0.1) is False


def test_fill_form_reports_overflow_warnings(tmp_path: Path) -> None:
    """fill_form 결과에 overflow_risk 플래그와 overflow_warnings 목록이 포함되고,
    넓은 값칸에 정상 배치되면 경고가 없어야 한다."""
    src = tmp_path / "f.hwpx"
    _make_form_hwpx(src, [("성명", ""), ("부서", "")])
    ad = load(src)
    r = ad.fill_form({"성명": "홍길동", "부서": "개발팀"})
    ad.close()
    assert "overflow_warnings" in r
    assert all("overflow_risk" in f for f in r["filled"])
    # 합성 2열 폼은 값칸이 충분히 넓어 오버플로 없음
    assert r["overflow_warnings"] == []


def _make_xlsx_form(path: Path) -> None:
    """병합 헤더 + 라벨-값 + 템플릿 키를 가진 xlsx 폼."""
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "신청서"
    ws["A1"] = "신청 정보"
    ws.merge_cells("A1:B1")
    ws["A2"] = "성명"
    ws["A3"] = "부서"
    ws["A4"] = "제목 {{title}}"
    wb.save(str(path))


def test_xlsx_inspect_fill_render_roundtrip(tmp_path: Path) -> None:
    """XlsxAdapter: 시트→표 인식, 병합, fill_form(base 상속), render, 영속."""
    src = tmp_path / "form.xlsx"
    _make_xlsx_form(src)
    ad = load(src)
    assert ad.format == "xlsx"
    t = ad.get_tables()[0]
    assert (t.rows, t.cols) == (4, 2)
    assert t.location == "신청서"
    assert [(m.anchor, m.span) for m in t.merges] == [((0, 0), (1, 2))]
    assert ad.get_placeholders() == ["title"]
    # 병합 anchor/non-anchor
    assert ad.get_cell(0, 0, 0).span == (1, 2)
    assert ad.get_cell(0, 0, 1).is_anchor is False
    # fill_form (base 구현이 자동 동작)
    r = ad.fill_form({"성명": "홍길동", "부서": "개발팀"})
    assert len(r["filled"]) == 2
    rr = ad.render_template({"title": "2026 보고서"})
    assert rr["used"] == ["title"]
    out = tmp_path / "out.xlsx"
    ad.save(out)
    ad.close()

    ad2 = load(out)
    assert ad2.get_cell(0, 1, 1).text == "홍길동"
    assert ad2.get_cell(0, 2, 1).text == "개발팀"
    assert "2026 보고서" in ad2.get_cell(0, 3, 0).text
    ad2.close()


def test_xlsx_value_typing(tmp_path: Path) -> None:
    """xlsx 셀 타입 처리: 날짜는 시간 제거, 금액은 숫자 보존, 전화/우편번호는 문자.

    회귀: 초기 구현은 set_cell 이 항상 문자로 써서 금액 셀이 텍스트가 되고
    (수식·합계 깨짐), 날짜를 '...00:00:00' 으로 표시했다.
    """
    import datetime
    from openpyxl import Workbook, load_workbook

    src = tmp_path / "t.xlsx"
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "날짜"
    ws["B1"] = datetime.date(2026, 6, 1)
    ws["A2"] = "금액"
    ws["A3"] = "전화"
    ws["A4"] = "우편"
    wb.save(str(src))

    ad = load(src)
    assert ad.get_cell(0, 0, 1).text == "2026-06-01"   # 시간 없음
    ad.set_cell(0, 1, 1, "3,000,000")    # 금액 → 숫자
    ad.set_cell(0, 2, 1, "010-1234-5678")  # 전화 → 문자
    ad.set_cell(0, 3, 1, "00100")          # 우편(선행0) → 문자
    ad.save(src)
    ad.close()

    wb2 = load_workbook(str(src))
    assert wb2.active["B2"].value == 3000000
    assert isinstance(wb2.active["B2"].value, int)
    assert wb2.active["B3"].value == "010-1234-5678"
    assert wb2.active["B4"].value == "00100"   # 선행 0 보존


def test_xlsx_complex_merges_formulas_roundtrip(tmp_path: Path) -> None:
    """복잡 xlsx: 가로+세로 병합, 수식, 통화서식, 다중시트가 편집 후에도 보존.

    set_cell 편집 시 다른 셀의 수식(=B*C, =SUM)·number_format·병합 영역·
    다른 시트가 그대로 유지돼야 한다 (openpyxl round-trip 정합성).
    """
    from openpyxl import Workbook, load_workbook

    src = tmp_path / "quote.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "견적서"
    ws["A1"] = "견적서"
    ws.merge_cells("A1:E1")               # 가로 병합 제목
    ws["A2"] = "공급자"
    ws.merge_cells("A2:A4")               # 세로 병합 섹션
    for r, (nm, qty, price) in enumerate(
            [("노트북", 2, 1500000), ("모니터", 3, 300000)], start=7):
        ws.cell(row=r, column=1, value=nm)
        ws.cell(row=r, column=2, value=qty)
        ws.cell(row=r, column=3, value=price)
        fcell = ws.cell(row=r, column=4, value=f"=B{r}*C{r}")
        fcell.number_format = '#,##0"원"'
    ws["D10"] = "=SUM(D7:D8)"
    wb.create_sheet("거래내역")["B1"] = 42
    wb.save(str(src))

    ad = load(src)
    t = ad.get_tables()[0]
    merges = {(m.anchor, m.span) for m in t.merges}
    assert ((0, 0), (1, 5)) in merges      # 가로 병합 제목
    assert ((1, 0), (3, 1)) in merges      # 세로 병합 섹션
    assert ad.get_cell(0, 0, 2).is_anchor is False     # 병합 non-anchor
    assert ad.get_cell(0, 6, 3).text == "=B7*C7"       # 수식 노출
    ad.set_cell(0, 6, 1, "10")             # 노트북 수량 편집
    ad.save(src)
    ad.close()

    wb2 = load_workbook(str(src))
    w = wb2["견적서"]
    assert w["B7"].value == 10                          # 숫자 보존
    assert w["D7"].value == "=B7*C7"                    # 수식 보존
    assert w["D10"].value == "=SUM(D7:D8)"
    assert w["D7"].number_format == '#,##0"원"'         # 서식 보존
    assert "A1:E1" in {str(m) for m in w.merged_cells.ranges}
    assert wb2["거래내역"]["B1"].value == 42            # 다른 시트 보존


def test_xlsx_merged_cell_write_rejected(tmp_path: Path) -> None:
    """병합 non-anchor 좌표 쓰기는 MergedCellWriteError (allow_merge_redirect로 우회)."""
    from document_adapter.base import MergedCellWriteError

    src = tmp_path / "f.xlsx"
    _make_xlsx_form(src)
    ad = load(src)
    with pytest.raises(MergedCellWriteError):
        ad.set_cell(0, 0, 1, "X")          # (0,1)은 A1:B1 병합의 non-anchor
    ad.set_cell(0, 0, 1, "X", allow_merge_redirect=True)  # anchor로 redirect
    assert ad.get_cell(0, 0, 0).text == "X"
    ad.close()


def test_xlsx_via_tools(tmp_path: Path) -> None:
    """MCP call_tool 경로(load 디스패치)로도 xlsx 가 동작."""
    from document_adapter.tools import call_tool

    src = tmp_path / "f.xlsx"
    _make_xlsx_form(src)
    insp = call_tool("inspect_document", {"path": str(src)})
    assert insp["format"] == "xlsx"
    assert insp["tables"][0]["location"] == "신청서"


def test_docx_header_footer_placeholders(tmp_path: Path) -> None:
    """docx get_placeholders 가 머리말/꼬리말의 {{key}} 도 포함해야 한다."""
    from docx import Document
    src = tmp_path / "hf.docx"
    d = Document()
    d.add_paragraph("본문 {{body}}")
    sec = d.sections[0]
    sec.header.paragraphs[0].text = "머리말 {{header_key}}"
    sec.footer.paragraphs[0].text = "꼬리말 {{footer_key}}"
    d.save(str(src))
    ad = load(src)
    ph = ad.get_placeholders()
    ad.close()
    assert {"body", "header_key", "footer_key"} <= set(ph)


def test_pptx_notes_placeholders_and_render(tmp_path: Path) -> None:
    """pptx 슬라이드 노트의 {{key}} 가 감지·렌더돼야 한다."""
    from pptx import Presentation
    src = tmp_path / "n.pptx"
    pr = Presentation()
    s = pr.slides.add_slide(pr.slide_layouts[6])
    s.notes_slide.notes_text_frame.text = "노트 {{note_key}}"
    pr.save(str(src))
    ad = load(src)
    assert "note_key" in ad.get_placeholders()
    ad.render_template({"note_key": "확인됨"})
    ad.save(src)
    ad.close()
    pr2 = Presentation(str(src))
    notes = [sl.notes_slide.notes_text_frame.text
             for sl in pr2.slides if sl.has_notes_slide]
    assert any("확인됨" in n for n in notes)


def test_get_cell_out_of_bounds_raises(tmp_path: Path) -> None:
    """경계를 벗어난 좌표는 CellOutOfBoundsError(IndexError 하위)."""
    from document_adapter.base import CellOutOfBoundsError

    src = tmp_path / "f.hwpx"
    _make_form_hwpx(src, [("a", "b")])
    ad = load(src)
    with pytest.raises(CellOutOfBoundsError):
        ad.get_cell(0, 99, 99)
    ad.close()
