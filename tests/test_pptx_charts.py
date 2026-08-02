"""PPTX 차트 기능 (v0.17): get_charts / set_chart_data / add_chart.

fixture 는 python-pptx 로 tmp_path 에 즉석 생성 (test_smoke.py 관례).
"""
from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from document_adapter import load
from document_adapter.base import NotImplementedForFormat
from document_adapter.tools import call_tool


CATS = ["1분기", "2분기", "3분기", "4분기"]
SALES = [120.0, 135.0, 150.0, 162.0]
PROFIT = [12.0, 18.0, 21.0, 25.0]


def _make_chart_pptx(path: Path) -> None:
    """column(2시리즈, 시리즈0 빨강) + pie(1시리즈) 2 슬라이드 덱."""
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    cd = CategoryChartData()
    cd.categories = CATS
    cd.add_series("매출", tuple(SALES))
    cd.add_series("영업이익", tuple(PROFIT))
    gf = s1.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1.5), Inches(8), Inches(5), cd,
    )
    ser0 = gf.chart.series[0]
    ser0.format.fill.solid()
    ser0.format.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)

    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    cd2 = CategoryChartData()
    cd2.categories = ["제품A", "제품B", "기타"]
    cd2.add_series("비중", (55.0, 30.0, 15.0))
    s2.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(2), Inches(1.5), Inches(6), Inches(5), cd2,
    )
    prs.save(path)


@pytest.fixture()
def chart_pptx(tmp_path: Path) -> Path:
    p = tmp_path / "charts.pptx"
    _make_chart_pptx(p)
    return p


def _chart_ids(path: Path) -> list[tuple[int, int]]:
    doc = load(path)
    try:
        return [(c.slide_index, c.shape_id) for c in doc.get_charts()]
    finally:
        doc.close()


# ---- get_charts ----

def test_get_charts_lists_all(chart_pptx: Path) -> None:
    doc = load(chart_pptx)
    try:
        charts = doc.get_charts()
    finally:
        doc.close()
    assert len(charts) == 2
    col = charts[0]
    assert col.slide_index == 1
    assert col.chart_type == "COLUMN_CLUSTERED"
    assert col.categories == CATS
    assert [s["name"] for s in col.series] == ["매출", "영업이익"]
    assert col.series[0]["values"] == SALES
    assert col.editable is True

    pie = charts[1]
    assert pie.slide_index == 2
    assert pie.chart_type == "PIE"
    assert pie.series[0]["values"] == [55.0, 30.0, 15.0]


def test_get_charts_slide_filter(chart_pptx: Path) -> None:
    doc = load(chart_pptx)
    try:
        assert len(doc.get_charts(slide_index=1)) == 1
        assert len(doc.get_charts(slide_index=2)) == 1
        assert doc.get_charts(slide_index=99) == []
    finally:
        doc.close()


def test_charts_invisible_to_tables_and_shapes(chart_pptx: Path) -> None:
    """차트는 get_tables/get_shapes 에 안 잡힌다 — get_charts 가 유일한 눈."""
    doc = load(chart_pptx)
    try:
        assert doc.get_tables() == []
        assert all(s.kind != "chart" for s in doc.get_shapes())
    finally:
        doc.close()


# ---- set_chart_data: set_points ----

def test_set_points_by_name(chart_pptx: Path) -> None:
    doc = load(chart_pptx)
    try:
        (s_idx, sh_id) = _chart_ids(chart_pptx)[0]
        result = doc.set_chart_data(
            s_idx, sh_id,
            set_points=[{"series": "매출", "category": "3분기", "value": 999}],
        )
        doc.save()
    finally:
        doc.close()
    assert result["before"]["series"][0]["values"] == SALES
    assert result["after"]["series"][0]["values"] == [120.0, 135.0, 999.0, 162.0]

    doc = load(chart_pptx)  # 재로드 검증
    try:
        chart = doc.get_charts()[0]
        assert chart.series[0]["values"] == [120.0, 135.0, 999.0, 162.0]
        assert chart.series[1]["values"] == PROFIT  # 다른 시리즈 불변
    finally:
        doc.close()


def test_set_points_by_index(chart_pptx: Path) -> None:
    doc = load(chart_pptx)
    try:
        (s_idx, sh_id) = _chart_ids(chart_pptx)[0]
        doc.set_chart_data(
            s_idx, sh_id,
            set_points=[{"series": 1, "category": 0, "value": "1,000"}],
        )
        doc.save()
    finally:
        doc.close()
    doc = load(chart_pptx)
    try:
        assert doc.get_charts()[0].series[1]["values"][0] == 1000.0
    finally:
        doc.close()


def test_set_points_unknown_series(chart_pptx: Path) -> None:
    doc = load(chart_pptx)
    try:
        (s_idx, sh_id) = _chart_ids(chart_pptx)[0]
        with pytest.raises(ValueError, match="매출"):
            # 없는 시리즈 이름 — 에러 메시지에 후보 목록 포함
            doc.set_chart_data(
                s_idx, sh_id,
                set_points=[{"series": "없는시리즈", "category": 0, "value": 1}],
            )
    finally:
        doc.close()


# ---- set_chart_data: 전체 교체 ----

def test_full_replace_grow_categories(chart_pptx: Path) -> None:
    new_cats = CATS + ["연간"]
    doc = load(chart_pptx)
    try:
        (s_idx, sh_id) = _chart_ids(chart_pptx)[0]
        doc.set_chart_data(
            s_idx, sh_id,
            categories=new_cats,
            series=[
                {"name": "매출", "values": [120, 135, 150, 162, 567]},
                {"name": "영업이익", "values": [12, 18, 21, 25, 76]},
            ],
        )
        doc.save()
    finally:
        doc.close()
    doc = load(chart_pptx)
    try:
        chart = doc.get_charts()[0]
        assert chart.categories == new_cats
        assert chart.series[0]["values"] == [120.0, 135.0, 150.0, 162.0, 567.0]
    finally:
        doc.close()


def test_full_replace_length_mismatch(chart_pptx: Path) -> None:
    doc = load(chart_pptx)
    try:
        (s_idx, sh_id) = _chart_ids(chart_pptx)[0]
        with pytest.raises(ValueError, match="values"):
            doc.set_chart_data(
                s_idx, sh_id,
                series=[{"name": "매출", "values": [1, 2]}],  # 카테고리는 4개
            )
        # 실패 시 문서 미변경
        assert doc.get_charts()[0].series[0]["values"] == SALES
    finally:
        doc.close()


def test_series_and_points_mutually_exclusive(chart_pptx: Path) -> None:
    doc = load(chart_pptx)
    try:
        (s_idx, sh_id) = _chart_ids(chart_pptx)[0]
        with pytest.raises(ValueError):
            doc.set_chart_data(
                s_idx, sh_id,
                series=[{"name": "매출", "values": [1, 2, 3, 4]}],
                set_points=[{"series": 0, "category": 0, "value": 1}],
            )
        with pytest.raises(ValueError):
            doc.set_chart_data(s_idx, sh_id)  # 아무것도 안 줌
    finally:
        doc.close()


def test_format_preserved_after_edit(chart_pptx: Path) -> None:
    """replace_data 후에도 시리즈 0 의 빨강(FF0000) solidFill 이 남아있어야 한다."""
    doc = load(chart_pptx)
    try:
        (s_idx, sh_id) = _chart_ids(chart_pptx)[0]
        doc.set_chart_data(
            s_idx, sh_id,
            set_points=[{"series": 0, "category": 0, "value": 1}],
        )
        doc.save()
    finally:
        doc.close()
    prs = Presentation(str(chart_pptx))
    chart_shape = next(
        s for s in prs.slides[0].shapes if getattr(s, "has_chart", False)
    )
    assert b"FF0000" in chart_shape.chart.part.blob


def test_title_only_change(chart_pptx: Path) -> None:
    doc = load(chart_pptx)
    try:
        (s_idx, sh_id) = _chart_ids(chart_pptx)[0]
        result = doc.set_chart_data(s_idx, sh_id, title="분기 실적")
        doc.save()
    finally:
        doc.close()
    assert result["title"] == "분기 실적"
    doc = load(chart_pptx)
    try:
        chart = doc.get_charts()[0]
        assert chart.title == "분기 실적"
        assert chart.series[0]["values"] == SALES  # 데이터 불변
    finally:
        doc.close()


# ---- add_chart ----

def test_add_chart_roundtrip(chart_pptx: Path) -> None:
    doc = load(chart_pptx)
    try:
        result = doc.add_chart(
            2, "line",
            categories=["1월", "2월", "3월"],
            series=[{"name": "온라인", "values": [10, 20, 30]},
                    {"name": "오프라인", "values": [5, 15, 25]}],
            title="월별 추이",
        )
        doc.save()
    finally:
        doc.close()
    assert result["chart_type"] == "LINE"

    doc = load(chart_pptx)
    try:
        charts = doc.get_charts(slide_index=2)
        added = next(c for c in charts if c.shape_id == result["shape_id"])
        assert added.chart_type == "LINE"
        assert added.categories == ["1월", "2월", "3월"]
        assert added.series[0]["values"] == [10.0, 20.0, 30.0]
        assert added.title == "월별 추이"
        # 추가된 차트도 곧장 편집 가능
        doc.set_chart_data(
            2, result["shape_id"],
            set_points=[{"series": "온라인", "category": "3월", "value": 99}],
        )
    finally:
        doc.close()


def test_add_chart_invalid_type(chart_pptx: Path) -> None:
    doc = load(chart_pptx)
    try:
        with pytest.raises(ValueError, match="column"):
            doc.add_chart(
                1, "3d_column",
                categories=["a"], series=[{"name": "s", "values": [1]}],
            )
    finally:
        doc.close()


def test_add_chart_default_placement(chart_pptx: Path) -> None:
    doc = load(chart_pptx)
    try:
        result = doc.add_chart(
            1, "bar",
            categories=["a", "b"],
            series=[{"name": "s", "values": [1, 2]}],
        )
        doc.save()
    finally:
        doc.close()
    prs = Presentation(str(chart_pptx))
    shape = next(
        s for s in prs.slides[0].shapes if s.shape_id == result["shape_id"]
    )
    assert shape.left >= 0 and shape.top >= 0
    assert shape.left + shape.width <= prs.slide_width
    assert shape.top + shape.height <= prs.slide_height


# ---- 포맷 경계 / 도구 dispatch ----

def test_non_pptx_rejected(tmp_path: Path) -> None:
    from docx import Document

    p = tmp_path / "doc.docx"
    Document().save(str(p))
    doc = load(p)
    try:
        assert doc.get_charts() == []
        with pytest.raises(NotImplementedForFormat):
            doc.set_chart_data(1, 1, title="x")
        with pytest.raises(NotImplementedForFormat):
            doc.add_chart(1, "column", categories=["a"],
                          series=[{"name": "s", "values": [1]}])
    finally:
        doc.close()


def test_tool_dispatch(chart_pptx: Path) -> None:
    listed = call_tool("get_charts", {"path": str(chart_pptx)})
    assert listed["chart_count"] == 2
    c = listed["charts"][0]

    edited = call_tool("set_chart_data", {
        "path": str(chart_pptx),
        "slide_index": c["slide_index"],
        "shape_id": c["shape_id"],
        "set_points": [{"series": "매출", "category": "1분기", "value": 777}],
    })
    assert "error" not in edited
    assert edited["after"]["series"][0]["values"][0] == 777.0

    added = call_tool("add_chart", {
        "path": str(chart_pptx),
        "slide_index": 1,
        "chart_type": "pie",
        "categories": ["x", "y"],
        "series": [{"name": "s", "values": [7, 3]}],
    })
    assert "error" not in added and added["chart_type"] == "PIE"


def test_tool_dispatch_not_implemented(tmp_path: Path) -> None:
    from docx import Document

    p = tmp_path / "doc.docx"
    Document().save(str(p))
    result = call_tool("set_chart_data", {
        "path": str(p), "slide_index": 1, "shape_id": 1, "title": "x",
    })
    assert result["error"] == "not_implemented"


def test_inspect_chart_summary(chart_pptx: Path, tmp_path: Path) -> None:
    result = call_tool("inspect_document", {"path": str(chart_pptx)})
    summary = result["chart_summary"]
    assert summary["chart_count"] == 2
    assert "set_chart_data" in summary["hint"]
    assert result["slide_count"] == 2

    # 차트 없는 pptx 엔 chart_summary 미포함
    p2 = tmp_path / "nochart.pptx"
    Presentation().save(str(p2))
    result2 = call_tool("inspect_document", {"path": str(p2)})
    assert "chart_summary" not in result2


def test_diff_chart_changes(chart_pptx: Path, tmp_path: Path) -> None:
    import shutil

    edited = tmp_path / "edited.pptx"
    shutil.copy2(chart_pptx, edited)
    call_tool("set_chart_data", {
        "path": str(edited), "slide_index": 1,
        "shape_id": _chart_ids(edited)[0][1],
        "set_points": [{"series": "매출", "category": "2분기", "value": 888}],
    })
    d = call_tool("diff_documents",
                  {"path_a": str(chart_pptx), "path_b": str(edited)})
    assert len(d["chart_changes"]) == 1
    ch = d["chart_changes"][0]
    assert ch["series"] == "매출" and ch["category"] == "2분기"
    assert ch["before"] == 135.0 and ch["after"] == 888.0

    # 차트 동일하면 chart_changes 키 자체가 없음 (기존 반환 형태 보존)
    d2 = call_tool("diff_documents",
                   {"path_a": str(chart_pptx), "path_b": str(chart_pptx)})
    assert "chart_changes" not in d2
