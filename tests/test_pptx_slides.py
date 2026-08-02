"""PPTX 슬라이드 기능 (v0.17): get_slides / duplicate_slide."""
from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from document_adapter import load
from document_adapter.base import NotImplementedForFormat
from document_adapter.tools import call_tool


def _make_deck(path: Path) -> None:
    """표지(텍스트) + 표 슬라이드 + 차트 슬라이드 3장 덱."""
    prs = Presentation()

    s1 = prs.slides.add_slide(prs.slide_layouts[0])  # 표지
    s1.shapes.title.text = "실적 보고"
    s1.placeholders[1].text = "작성자: 테스터"

    s2 = prs.slides.add_slide(prs.slide_layouts[5])  # 표
    s2.shapes.title.text = "분기별 표"
    tbl = s2.shapes.add_table(
        2, 2, Inches(1), Inches(2), Inches(6), Inches(2)
    ).table
    tbl.cell(0, 0).text = "구분"
    tbl.cell(0, 1).text = "값"
    tbl.cell(1, 0).text = "매출"
    tbl.cell(1, 1).text = "100"

    s3 = prs.slides.add_slide(prs.slide_layouts[5])  # 차트
    s3.shapes.title.text = "차트 페이지"
    cd = CategoryChartData()
    cd.categories = ["1분기", "2분기"]
    cd.add_series("매출", (100.0, 200.0))
    s3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(2), Inches(8), Inches(4), cd,
    )
    prs.save(path)


@pytest.fixture()
def deck(tmp_path: Path) -> Path:
    p = tmp_path / "deck.pptx"
    _make_deck(p)
    return p


# ---- get_slides ----

def test_get_slides_overview(deck: Path) -> None:
    doc = load(deck)
    try:
        slides = doc.get_slides()
    finally:
        doc.close()
    assert len(slides) == 3
    assert slides[0].slide_index == 1
    assert slides[0].title == "실적 보고"
    assert slides[0].table_count == 0 and slides[0].chart_count == 0

    assert slides[1].title == "분기별 표"
    assert slides[1].table_count == 1

    assert slides[2].title == "차트 페이지"
    assert slides[2].chart_count == 1
    assert all(s.layout_name for s in slides)


def test_get_slides_non_pptx(tmp_path: Path) -> None:
    from docx import Document

    p = tmp_path / "doc.docx"
    Document().save(str(p))
    doc = load(p)
    try:
        assert doc.get_slides() == []
        with pytest.raises(NotImplementedForFormat):
            doc.duplicate_slide(1)
    finally:
        doc.close()


# ---- duplicate_slide ----

def test_duplicate_table_slide_append(deck: Path) -> None:
    doc = load(deck)
    try:
        result = doc.duplicate_slide(2)
        doc.save()
    finally:
        doc.close()
    assert result["new_slide_index"] == 4
    assert result["slide_count"] == 4
    assert len(result["tables"]) == 1
    assert result["tables"][0]["rows"] == 2
    # preview 로 행 라벨이 보여야 함 — LLM 이 행 매핑을 추측하지 않게
    assert result["tables"][0]["preview"][0][0] == "구분"
    assert result["tables"][0]["preview"][1][0] == "매출"
    assert "warning" not in result  # 맨 뒤 추가는 index 시프트 없음

    doc = load(deck)
    try:
        slides = doc.get_slides()
        assert len(slides) == 4
        assert slides[3].title == "분기별 표"
        assert slides[3].table_count == 1
        # 복제된 표 내용 확인 (전역 index 는 반환값 기준)
        t_idx = result["tables"][0]["table_index"]
        assert doc.get_cell(t_idx, 1, 1).text == "100"
    finally:
        doc.close()


def test_duplicate_at_position(deck: Path) -> None:
    doc = load(deck)
    try:
        result = doc.duplicate_slide(2, at=2)
        doc.save()
    finally:
        doc.close()
    assert result["new_slide_index"] == 2
    assert "table_index" in result["warning"]

    doc = load(deck)
    try:
        slides = doc.get_slides()
        assert [s.title for s in slides] == [
            "실적 보고", "분기별 표", "분기별 표", "차트 페이지",
        ]
    finally:
        doc.close()


def test_table_index_shift_reported(deck: Path) -> None:
    """중간 삽입 후 반환된 table_index 가 삽입 후 좌표계와 일치해야 한다."""
    doc = load(deck)
    try:
        result = doc.duplicate_slide(2, at=1)  # 맨 앞 삽입 → 기존 표는 뒤로
        # 반환 좌표로 바로 편집 가능해야 함
        t_idx = result["tables"][0]["table_index"]
        assert t_idx == 0  # 맨 앞 슬라이드의 표가 전역 0번
        doc.set_cell(t_idx, 1, 1, "999")
        doc.save()
    finally:
        doc.close()

    doc = load(deck)
    try:
        # 복제본(slide 1)만 변경, 원본(이제 slide 3)은 불변
        assert doc.get_cell(0, 1, 1).text == "999"
        assert doc.get_cell(1, 1, 1).text == "100"
    finally:
        doc.close()


def test_duplicate_chart_slide_parts_cloned(deck: Path) -> None:
    doc = load(deck)
    try:
        result = doc.duplicate_slide(3)
        doc.save()
    finally:
        doc.close()
    assert len(result["charts"]) == 1

    prs = Presentation(str(deck))
    def chart_part(slide):
        return next(
            s.chart.part.partname for s in slide.shapes
            if getattr(s, "has_chart", False)
        )
    p_orig = chart_part(prs.slides[2])
    p_dup = chart_part(prs.slides[3])
    assert p_orig != p_dup  # chart part 독립 복제


def test_edit_duplicate_leaves_original(deck: Path) -> None:
    doc = load(deck)
    try:
        result = doc.duplicate_slide(3)
        doc.set_chart_data(
            result["new_slide_index"], result["charts"][0]["shape_id"],
            set_points=[{"series": "매출", "category": "1분기", "value": 777}],
        )
        doc.save()
    finally:
        doc.close()

    doc = load(deck)
    try:
        orig = doc.get_charts(slide_index=3)[0]
        dup = doc.get_charts(slide_index=4)[0]
        assert orig.series[0]["values"] == [100.0, 200.0]  # 원본 불변
        assert dup.series[0]["values"] == [777.0, 200.0]
    finally:
        doc.close()


def test_duplicate_invalid_index(deck: Path) -> None:
    doc = load(deck)
    try:
        with pytest.raises(ValueError):
            doc.duplicate_slide(99)
        with pytest.raises(ValueError):
            doc.duplicate_slide(1, at=99)
        with pytest.raises(ValueError):
            doc.duplicate_slide(0)
    finally:
        doc.close()


def test_roundtrip_after_duplicate(deck: Path) -> None:
    """복제 후 전역 인덱스 일관성: tables/shapes/charts 재조회가 서로 맞아야 함."""
    doc = load(deck)
    try:
        doc.duplicate_slide(2, at=2)
        doc.save()
    finally:
        doc.close()

    doc = load(deck)
    try:
        tables = doc.get_tables()
        assert len(tables) == 2
        assert tables[0].location == "slide 2"
        assert tables[1].location == "slide 3"
        assert [t.index for t in tables] == [0, 1]
        charts = doc.get_charts()
        assert len(charts) == 1 and charts[0].slide_index == 4
    finally:
        doc.close()


# ---- 도구 dispatch ----

def test_tool_dispatch_slides(deck: Path) -> None:
    listed = call_tool("get_slides", {"path": str(deck)})
    assert listed["slide_count"] == 3
    assert listed["slides"][1]["table_count"] == 1

    dup = call_tool("duplicate_slide", {
        "path": str(deck), "source_slide_index": 2,
    })
    assert "error" not in dup
    assert dup["new_slide_index"] == 4

    # 반환 좌표로 즉시 편집
    t_idx = dup["tables"][0]["table_index"]
    edited = call_tool("set_cell", {
        "path": str(deck), "table_index": t_idx,
        "row": 1, "col": 1, "value": "888",
    })
    assert edited["previous_value"] == "100"


def test_tool_dispatch_duplicate_output_path(deck: Path, tmp_path: Path) -> None:
    out = tmp_path / "copy.pptx"
    result = call_tool("duplicate_slide", {
        "path": str(deck), "source_slide_index": 1,
        "output_path": str(out),
    })
    assert "error" not in result
    # 원본은 3장 유지, output 은 4장
    assert call_tool("get_slides", {"path": str(deck)})["slide_count"] == 3
    assert call_tool("get_slides", {"path": str(out)})["slide_count"] == 4
