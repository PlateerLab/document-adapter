"""블록 IR → PPTX 결정적 렌더러 (LLM 없음).

markdown 서브셋을 슬라이드 덱으로 변환한다. docx_writer 와 같은 파서
(markdown_parser)를 공유하고, 슬라이드 분할 규칙만 얹는다:

- ``---`` (hr) 블록 = 슬라이드 구분자 (Marp/reveal.js 관례)
- 레벨 1~2 헤딩은 새 슬라이드 시작 (현재 슬라이드에 내용이 이미 있으면)
- 슬라이드의 첫 헤딩 = 슬라이드 제목, 나머지 블록 = 본문

렌더 원칙:
- python-pptx **기본 템플릿의 내장 레이아웃만** 사용
  (1 = Title and Content, 5 = Title Only). 커스텀 마스터 정의 금지.
- 표가 없는 슬라이드 → 레이아웃 1: 본문 placeholder 에 불릿/문단.
- 표가 있는 슬라이드 → 레이아웃 5: 본문 텍스트박스(있을 때) + 표를
  고정 좌표에 순서대로 배치 (겹침 없는 결정적 배치).
- 생성 직후 PptxAdapter.load() roundtrip 이 성립해야 한다
  (get_shapes/set_shape_text/set_cell 로 이어 편집).
"""
from __future__ import annotations

import io

from pptx import Presentation
from pptx.util import Emu, Pt

from .docx_writer import base_font_for_lang
from .markdown_parser import Block, Span, parse_markdown

__all__ = ["pptx_from_markdown"]

_CODE_FONT = "D2Coding"

# 기본 템플릿(13.33x7.5in) 기준 고정 배치 좌표 (EMU).
_EMU_PER_IN = 914400
_SLIDE_W = int(13.333 * _EMU_PER_IN)
_MARGIN_X = int(0.6 * _EMU_PER_IN)
_TITLE_BOTTOM = int(1.35 * _EMU_PER_IN)   # 레이아웃 5 제목 아래 y
_BODY_H = int(1.8 * _EMU_PER_IN)          # 표 동반 슬라이드의 본문 텍스트박스 높이
_TABLE_ROW_H = int(0.38 * _EMU_PER_IN)
_CONTENT_W = _SLIDE_W - 2 * _MARGIN_X


def _split_slides(blocks: list[Block]) -> list[list[Block]]:
    """블록 리스트 → 슬라이드별 블록 리스트.

    hr = 명시적 구분자. 레벨 1~2 헤딩은 (현재 슬라이드에 내용이 있으면)
    새 슬라이드를 시작한다 — ``# 제목`` 만 반복하는 문서도 자연히 덱이 된다.
    """
    slides: list[list[Block]] = []
    current: list[Block] = []
    for block in blocks:
        if block.kind == "hr":
            if current:
                slides.append(current)
                current = []
            continue
        if block.kind == "heading" and block.level <= 2 and current:
            slides.append(current)
            current = []
        current.append(block)
    if current:
        slides.append(current)
    return slides


def _spans_into_paragraph(paragraph, spans: tuple[Span, ...], font: str,
                          size: Pt | None = None) -> None:
    for span in spans:
        run = paragraph.add_run()
        run.text = span.text
        run.font.name = _CODE_FONT if span.code else font
        if size is not None:
            run.font.size = size
        if span.bold:
            run.font.bold = True
        if span.italic:
            run.font.italic = True


def _fill_text_frame(tf, body: list[Block], font: str) -> None:
    """본문 블록들을 text frame 문단으로. 첫 문단은 기존 것 재사용."""
    tf.word_wrap = True
    first = True
    numbered_count = 0
    for block in body:
        if block.kind == "code":
            for line in block.lines:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                run = p.add_run()
                run.text = line
                run.font.name = _CODE_FONT
                run.font.size = Pt(12)
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        spans = block.spans
        if block.kind == "numbered":
            # placeholder 자동 번호 대신 결정적 텍스트 프리픽스.
            numbered_count += 1
            spans = (Span(f"{numbered_count}. "),) + spans
        else:
            numbered_count = 0
        if block.kind == "quote":
            p.level = 1
        if block.kind == "heading":
            # 슬라이드 내부 소제목 (레벨 3+) — 굵게.
            spans = tuple(
                Span(s.text, bold=True, italic=s.italic, code=s.code) for s in spans
            )
        _spans_into_paragraph(p, spans, font)


def _add_table_shape(slide, block: Block, top: int, font: str) -> int:
    """표를 고정 좌표에 추가하고 다음 top(EMU)을 반환."""
    n_rows = len(block.rows)
    n_cols = len(block.rows[0]) if block.rows else 0
    if n_rows == 0 or n_cols == 0:
        return top
    height = _TABLE_ROW_H * n_rows
    shape = slide.shapes.add_table(
        n_rows, n_cols, Emu(_MARGIN_X), Emu(top), Emu(_CONTENT_W), Emu(height)
    )
    table = shape.table
    for r, row in enumerate(block.rows):
        for c, cell_spans in enumerate(row):
            cell = table.cell(r, c)
            para = cell.text_frame.paragraphs[0]
            spans = cell_spans
            if r == 0:  # 헤더행 굵게
                spans = tuple(
                    Span(s.text, bold=True, italic=s.italic, code=s.code)
                    for s in spans
                )
            _spans_into_paragraph(para, spans, font, size=Pt(14))
    return top + height + int(0.25 * _EMU_PER_IN)


def pptx_from_markdown(
    markdown: str, *, lang: str | None = "ko", base_font: str | None = None
) -> bytes:
    """markdown 서브셋 → 슬라이드 덱 .pptx bytes.

    Raises:
        ValueError: 본문이 비어 있을 때 (이중어 메시지 — 호출 레이어의
            재시도 계약이 이 메시지를 LLM 리마인더로 사용한다).
    """
    if not markdown or not markdown.strip():
        raise ValueError(
            "empty document body — provide markdown content. "
            "문서 본문이 비어 있습니다 — markdown 내용을 작성하세요."
        )

    blocks = parse_markdown(markdown)
    slides_blocks = _split_slides(blocks)
    if not slides_blocks:
        raise ValueError(
            "markdown produced no renderable slides. "
            "markdown 에서 렌더 가능한 슬라이드를 찾지 못했습니다."
        )

    font = base_font or base_font_for_lang(lang)
    prs = Presentation()
    layout_content = prs.slide_layouts[1]  # Title and Content
    layout_title_only = prs.slide_layouts[5]  # Title Only
    prs.slide_width = Emu(_SLIDE_W)
    prs.slide_height = Emu(int(7.5 * _EMU_PER_IN))

    for slide_blocks in slides_blocks:
        title_spans: tuple[Span, ...] = ()
        body: list[Block] = []
        tables: list[Block] = []
        for block in slide_blocks:
            if block.kind == "heading" and not title_spans:
                title_spans = block.spans
            elif block.kind == "table":
                tables.append(block)
            else:
                body.append(block)

        slide = prs.slides.add_slide(
            layout_title_only if tables else layout_content
        )
        title_shape = slide.shapes.title
        if title_shape is not None:
            para = title_shape.text_frame.paragraphs[0]
            _spans_into_paragraph(
                para, title_spans or (Span(""),), font, size=Pt(30)
            )

        if tables:
            top = _TITLE_BOTTOM
            if body:
                tb = slide.shapes.add_textbox(
                    Emu(_MARGIN_X), Emu(top), Emu(_CONTENT_W), Emu(_BODY_H)
                )
                _fill_text_frame(tb.text_frame, body, font)
                top += _BODY_H + int(0.15 * _EMU_PER_IN)
            for table_block in tables:
                top = _add_table_shape(slide, table_block, top, font)
        elif body:
            # 레이아웃 1 의 콘텐츠 placeholder (idx 1) 에 본문.
            placeholder = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    placeholder = ph
                    break
            if placeholder is not None:
                _fill_text_frame(placeholder.text_frame, body, font)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
