"""공개 합성 양식 fixture 빌더.

LLM 평가에는 양식 문서가 필요한데, 실제 양식(``tests/fixtures/*/real/``)은
개인·회사 문서라 커밋할 수 없다. 여기서는 동일한 *구조 패턴*을 코드로
재현한 공개 fixture 를 만든다 — 누구나 재현 가능하고 CI 에서도 쓸 수 있다.

패턴:
  - label-right : 라벨(좌) | 값(우)  — 가장 흔한 양식. fill_form auto.
  - dup-section : 같은 라벨이 두 섹션에 — dot-path 로 구분해야 함.
  - filled-sample : 값 칸에 예시값이 미리 들어있음 — direction="right" 필요.

각 포맷별 빌더는 docx/pptx(런타임 의존) 와 hwpx(dev 의존 python-hwpx)를
**지연 import** 한다 — 이 모듈 import 만으로 무거운 의존을 끌어오지 않도록.
"""
from __future__ import annotations

from pathlib import Path


# ---- HWPX -----------------------------------------------------------------

def _hwpx_form(path: Path, rows: list[tuple[str, str]]) -> None:
    """라벨-값 2열 표 HWPX. rows = [(label, value), ...]."""
    from hwpx.document import HwpxDocument

    doc = HwpxDocument.new()
    doc.add_paragraph("")
    doc.add_table(len(rows), 2)
    doc.save_to_path(path)

    d = HwpxDocument.open(path)
    try:
        sec = d.sections[0]
        tbl = next(p.tables[0] for p in sec.paragraphs if p.tables)
        for i, (label, value) in enumerate(rows):
            tbl.rows[i].cells[0].text = label
            tbl.rows[i].cells[1].text = value
        d.save_to_path(path)
    finally:
        d.close()


def hwpx_blank_form(path: Path) -> None:
    """빈 값칸 양식 (label-right, auto)."""
    _hwpx_form(path, [
        ("접수번호", ""),
        ("성명", ""),
        ("주소", ""),
    ])


def _hwpx_grid(
    path: Path,
    cells: list[list[str]],
    merges: list[tuple[int, int, int, int]] | None = None,
) -> None:
    """임의 rows×cols HWPX 표 + 선택적 병합.

    cells: 행별 셀 텍스트 2차원 리스트.
    merges: (row, col, row_span, col_span) 앵커 병합 목록. 병합에 덮이는
            non-anchor 셀은 width/height 0 으로 비활성화한다(HWPX 병합 관례).
    """
    from hwpx.document import HwpxDocument

    rows, cols = len(cells), len(cells[0])
    doc = HwpxDocument.new()
    doc.add_paragraph("")
    doc.add_table(rows, cols)
    doc.save_to_path(path)

    d = HwpxDocument.open(path)
    try:
        sec = d.sections[0]
        tbl = next(p.tables[0] for p in sec.paragraphs if p.tables)
        covered: set[tuple[int, int]] = set()
        for (r, c, rs, cs) in (merges or []):
            anchor = tbl.rows[r].cells[c]
            anchor.set_span(row_span=rs, col_span=cs)
            for rr in range(r, r + rs):
                for cc in range(c, c + cs):
                    if (rr, cc) != (r, c):
                        tbl.rows[rr].cells[cc].set_size(width=0, height=0)
                        tbl.rows[rr].cells[cc].text = ""
                        covered.add((rr, cc))
        for r in range(rows):
            for c in range(cols):
                if (r, c) in covered:
                    continue
                tbl.rows[r].cells[c].text = cells[r][c]
        d.save_to_path(path)
    finally:
        d.close()


def hwpx_header_below_form(path: Path) -> None:
    """헤더행 라벨 / 값은 아래 행 — 명부·집계표류 (direction below)."""
    _hwpx_grid(path, [
        ["성명", "부서", "직급"],
        ["", "", ""],
    ])


def hwpx_wide_grid_form(path: Path) -> None:
    """인적사항 4열 그리드 (이력서·신상카드류): 라벨|값|라벨|값."""
    _hwpx_grid(path, [
        ["성명", "", "생년월일", ""],
        ["주소", "", "연락처", ""],
    ])


def hwpx_merged_section_form(path: Path) -> None:
    """섹션 헤더가 colspan 병합된 신청서류: 헤더 아래 라벨-값 행."""
    _hwpx_grid(
        path,
        [
            ["신청인 정보", ""],
            ["성명", ""],
            ["연락처", ""],
        ],
        merges=[(0, 0, 1, 2)],   # 0행을 colspan=2 로 병합한 섹션 헤더
    )


def hwpx_three_section_form(path: Path) -> None:
    """신청인·대리인·보증인 3개 섹션에 동일 라벨(성명/연락처) — dot-path x3."""
    _hwpx_grid(path, [
        ["신청인", ""], ["성명", ""], ["연락처", ""],
        ["대리인", ""], ["성명", ""], ["연락처", ""],
        ["보증인", ""], ["성명", ""], ["연락처", ""],
    ])


def hwpx_dup_section_form(path: Path) -> None:
    """같은 '금액' 라벨이 두 섹션에 — dot-path 구분 필요.

    섹션 헤더(col 0)는 _candidate_context_labels 가 위로 훑어 찾으므로,
    각 섹션 첫 행에 헤더 라벨을 둔다.
    """
    _hwpx_form(path, [
        ("피해자 정보", ""),
        ("금액", ""),
        ("지급정지요청계좌", ""),
        ("금액", ""),
    ])


# ---- DOCX -----------------------------------------------------------------

def docx_blank_form(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_paragraph("직원 정보 양식")
    rows = [("성명", ""), ("부서", ""), ("입사일", "")]
    table = doc.add_table(rows=len(rows), cols=2)
    for i, (label, value) in enumerate(rows):
        table.cell(i, 0).text = label
        table.cell(i, 1).text = value
    doc.save(str(path))


# ---- PPTX -----------------------------------------------------------------

def pptx_filled_sample_form(path: Path) -> None:
    """값칸에 예시값이 들어있는 PPTX 표 — direction='right' 로 교체해야 함."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    rows, cols = 3, 2
    gtbl = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(0.5), Inches(9), Inches(3)
    ).table
    data = [("보고일자", "0000-00-00"), ("작성자", "예시"), ("담당부서", "예시팀")]
    for i, (label, sample) in enumerate(data):
        gtbl.cell(i, 0).text = label
        gtbl.cell(i, 1).text = sample
    prs.save(str(path))
