"""PPTX 어댑터: python-pptx + 자체 {{key}} 치환 엔진.

표 구조:
- python-pptx는 ``cell.is_merge_origin`` / ``cell.is_spanned`` /
  ``cell.span_height`` / ``cell.span_width`` 로 병합 정보를 직접 노출.
- PPTX는 중첩 테이블이 없음 (셀은 text_frame만 보유).
"""
from __future__ import annotations

import datetime
import re
import warnings
from copy import deepcopy
from pathlib import Path
from typing import Any, Iterator

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Emu

from .base import (
    CellContent,
    CellOutOfBoundsError,
    ChartInfo,
    DocumentAdapter,
    MergeInfo,
    MergedCellWriteError,
    NotImplementedForFormat,
    ShapeInfo,
    SlideInfo,
    TableIndexError,
    TableSchema,
    _has_template,
)

TAG_PATTERN = re.compile(r"\{\{\s*(\w+)\s*\}\}")
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

# OOXML EMU (English Metric Unit) → cm: 1 cm = 360000 EMU
_EMU_PER_CM = 360000

# LLM 친화 문자열 → XL_CHART_TYPE (v1: category 계열만 — scatter/bubble/3D 는
# 데이터 모델이 달라 편집 계약이 흔들리므로 제외. base.ChartInfo docstring 참조).
_CHART_TYPE_MAP: dict[str, Any] = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
    "area_stacked": XL_CHART_TYPE.AREA_STACKED,
    "radar": XL_CHART_TYPE.RADAR,
}

# add_chart 기본 배치 (cm) — 제목 아래 전폭 (pptx_writer 의 결정적 배치 원칙)
_CHART_DEFAULT_X_CM = 1.5
_CHART_DEFAULT_Y_CM = 3.5


def _iter_shapes_recursive(shapes: Any) -> Iterator[Any]:
    """그룹 shape 안까지 재귀 순회 (차트/집계는 그룹 내부도 봐야 한다)."""
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes_recursive(shape.shapes)


def _to_number(v: Any) -> float | int | None:
    """차트 값 정규화. None 허용, 문자열은 콤마 제거 후 float 시도."""
    if v is None:
        return None
    if isinstance(v, bool):
        raise ValueError(
            f"chart value must be a number, got bool {v!r}. "
            f"차트 값은 숫자여야 합니다 (bool 불가)."
        )
    if isinstance(v, (int, float)):
        return v
    if isinstance(v, str):
        s = v.replace(",", "").strip()
        try:
            return float(s)
        except ValueError:
            pass
    raise ValueError(
        f"not a number: {v!r} — remove units/commas and pass a number. "
        f"숫자가 아닌 값입니다: {v!r} — 단위/콤마를 제거하고 숫자로 전달하세요."
    )


def _resolve_index(key: Any, names: list[str], what: str) -> int:
    """이름 우선, 정수 인덱스 허용으로 시리즈/카테고리를 좌표로 해소."""
    if isinstance(key, bool) or key is None:
        raise ValueError(f"{what} must be a name or 0-based index, got {key!r}")
    if isinstance(key, int):
        if 0 <= key < len(names):
            return key
        raise ValueError(
            f"{what} index {key} out of range 0..{len(names) - 1}"
        )
    k = str(key).strip()
    matches = [i for i, n in enumerate(names) if (n or "").strip() == k]
    if len(matches) == 1:
        return matches[0]
    if not matches:
        raise ValueError(
            f"{what} not found: {key!r}. candidates: {names}. "
            f"{what} 이름을 찾지 못했습니다 — 후보 중에서 고르거나 인덱스를 쓰세요."
        )
    raise ValueError(
        f"{what} name {key!r} is ambiguous (indices {matches}); "
        f"use an integer index instead."
    )


def _validate_series(series: Any, categories: list[str]) -> list[dict[str, Any]]:
    """series 인자 검증/정규화: [{"name", "values"}], 길이 == 카테고리 수."""
    if not isinstance(series, (list, tuple)) or not series:
        raise ValueError(
            "series must be a non-empty list of {name, values}. "
            "series 는 {name, values} dict 의 비어있지 않은 리스트여야 합니다."
        )
    out: list[dict[str, Any]] = []
    for i, s in enumerate(series):
        if not isinstance(s, dict) or "values" not in s:
            raise ValueError(
                f"series[{i}] must be a dict with 'name' and 'values'."
            )
        name = str(s.get("name") or f"Series {i + 1}")
        vals = s["values"]
        if not isinstance(vals, (list, tuple)):
            raise ValueError(f"series[{i}].values must be a list of numbers")
        if len(vals) != len(categories):
            raise ValueError(
                f"series[{i}] ({name!r}) has {len(vals)} values but there are "
                f"{len(categories)} categories. "
                f"값 개수({len(vals)})가 카테고리 수({len(categories)})와 다릅니다."
            )
        out.append({"name": name, "values": [_to_number(v) for v in vals]})
    return out


def _clone_part(part: Any) -> Any:
    """OPC part 를 같은 클래스/blob 으로 새 partname 에 복제.

    관계(rels)는 원본과 같은 rId 오름차순으로 재생성한다 — 빈 rels 에서
    시작하므로 동일한 rId 가 부여되어, blob 안의 ``r:id`` 참조를 재작성하지
    않아도 유효하다. 차트 part 의 내장 xlsx(``RT.PACKAGE``)는 재귀 복제해
    복제본 차트의 데이터 편집이 원본을 오염시키지 않게 한다.
    """
    pkg = part.package
    partname = str(part.partname)
    base, ext = partname.rsplit(".", 1)
    tmpl = base.rstrip("0123456789") + "%d." + ext
    new_name = pkg.next_partname(tmpl)
    new_part = type(part).load(new_name, part.content_type, pkg, part.blob)
    for rid in sorted(part.rels, key=lambda r: int(r[3:])):
        rel = part.rels[rid]
        if rel.is_external:
            new_rid = new_part.rels.get_or_add_ext_rel(
                rel.reltype, rel.target_ref
            )
        else:
            target = rel.target_part
            if rel.reltype == RT.PACKAGE:
                target = _clone_part(target)
            new_rid = new_part.relate_to(target, rel.reltype)
        if new_rid != rid:
            # rId 재현 가정이 깨진 문서 — 조용한 참조 오염 대신 명시적 거부.
            raise NotImplementedForFormat(
                f"cannot safely clone part {partname}: relationship id "
                f"mismatch ({new_rid} != {rid}). "
                f"이 문서의 관계 구조는 안전한 복제를 지원하지 않습니다."
            )
    return new_part


def _emu_to_cm(emu: Any) -> float | None:
    """EMU 값을 cm 1자리 반올림. None 또는 0 은 None."""
    if emu is None:
        return None
    try:
        v = int(emu)
    except (TypeError, ValueError):
        return None
    if v <= 0:
        return None
    return round(v / _EMU_PER_CM, 1)


class PptxAdapter(DocumentAdapter):
    format = "pptx"

    def _open(self) -> None:
        self._prs = Presentation(str(self.path))

    def save(self, path: Path | str | None = None) -> Path:
        target = Path(path) if path else self.path
        self._prs.save(str(target))
        self.path = target
        return target

    # ---- helpers ----
    def _iter_tables(self) -> Iterator[tuple[int, int, Any]]:
        """(global_index, slide_number_1based, table) 순회."""
        g_idx = 0
        for s_idx, slide in enumerate(self._prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_table:
                    yield g_idx, s_idx, shape.table
                    g_idx += 1

    def _iter_text_frames(self) -> Iterator[Any]:
        for slide in self._prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    yield shape.text_frame
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            yield cell.text_frame
            # 슬라이드 노트의 {{key}} 도 포함 (get_placeholders·render 공통)
            if slide.has_notes_slide:
                yield slide.notes_slide.notes_text_frame

    @staticmethod
    def _dimensions(table) -> tuple[int, int]:
        n_rows = len(list(table.rows))
        n_cols = len(list(table.columns))
        return n_rows, n_cols

    def _resolve_anchor_cell(
        self, table, row: int, col: int, *, allow_merge_redirect: bool
    ) -> tuple[Any, tuple[int, int], tuple[int, int], bool]:
        """(cell, anchor, span, is_anchor) 반환.

        non-anchor 좌표이고 allow_merge_redirect=False면 MergedCellWriteError.
        True면 anchor cell로 리디렉트하고 경고.
        """
        n_rows, n_cols = self._dimensions(table)
        if row < 0 or col < 0 or row >= n_rows or col >= n_cols:
            raise CellOutOfBoundsError(
                f"cell ({row},{col}) out of bounds ({n_rows}x{n_cols})"
            )
        cell = table.cell(row, col)

        # merge anchor 좌표 계산
        is_anchor = bool(getattr(cell, "is_merge_origin", False)) or not bool(
            getattr(cell, "is_spanned", False)
        )
        if cell.is_merge_origin:
            span = (cell.span_height, cell.span_width)
            anchor = (row, col)
        elif cell.is_spanned:
            # anchor는 앞쪽 어딘가. 탐색으로 origin 찾기
            anchor = self._find_merge_origin(table, row, col)
            origin_cell = table.cell(*anchor)
            span = (origin_cell.span_height, origin_cell.span_width)
            is_anchor = False
            if not allow_merge_redirect:
                raise MergedCellWriteError(
                    f"cell ({row},{col}) is part of a merged region anchored at "
                    f"({anchor[0]},{anchor[1]}) span={span}. "
                    f"Write to the anchor coordinate, or pass "
                    f"allow_merge_redirect=True."
                )
            warnings.warn(
                f"write to ({row},{col}) redirected to merge anchor "
                f"({anchor[0]},{anchor[1]})",
                stacklevel=3,
            )
            cell = origin_cell
        else:
            span = (1, 1)
            anchor = (row, col)

        return cell, anchor, span, is_anchor

    @staticmethod
    def _find_merge_origin(table, row: int, col: int) -> tuple[int, int]:
        """is_spanned 셀로부터 병합 origin 좌표를 역추적.

        python-pptx가 origin 좌표 자체를 직접 노출하지 않아 앞쪽 row/col을 훑어
        이 (row,col)을 포함하는 origin을 찾는다. 테이블 크기가 작을 때는 충분.
        """
        for r in range(row, -1, -1):
            for c in range(col, -1, -1):
                candidate = table.cell(r, c)
                if not candidate.is_merge_origin:
                    continue
                if (r + candidate.span_height > row) and (
                    c + candidate.span_width > col
                ):
                    return (r, c)
        # fallback
        return (row, col)

    # ---- inspection ----
    def get_placeholders(self) -> list[str]:
        keys: set[str] = set()
        for tf in self._iter_text_frames():
            keys.update(TAG_PATTERN.findall(tf.text))
        return sorted(keys)

    def get_tables(self, min_rows: int = 1, min_cols: int = 1,
                   preview_rows: int = 4, max_cell_len: int = 40) -> list[TableSchema]:
        schemas: list[TableSchema] = []
        for g_idx, s_idx, table in self._iter_tables():
            n_rows, n_cols = self._dimensions(table)
            if n_rows < min_rows or n_cols < min_cols:
                continue

            visible_rows = min(n_rows, preview_rows)
            preview: list[list[str | None]] = [
                [None for _ in range(n_cols)] for _ in range(visible_rows)
            ]
            merges: list[MergeInfo] = []

            for r in range(n_rows):
                for c in range(n_cols):
                    cell = table.cell(r, c)
                    if cell.is_spanned:
                        continue  # non-anchor, preview stays None
                    # anchor (merge origin or standalone cell)
                    if r < visible_rows:
                        text = (cell.text or "").strip()
                        preview[r][c] = text[:max_cell_len]
                    if cell.is_merge_origin:
                        span = (cell.span_height, cell.span_width)
                        if span != (1, 1):
                            merges.append(MergeInfo(anchor=(r, c), span=span))

            # 셀 크기 힌트 (EMU → cm). LLM 이 오버플로 위험 셀을 판단하는 데 사용.
            col_widths = [
                _emu_to_cm(getattr(col, "width", None)) for col in table.columns
            ]
            row_heights = [
                _emu_to_cm(getattr(row, "height", None)) for row in table.rows
            ]
            col_widths_out = col_widths if any(v is not None for v in col_widths) else None
            row_heights_out = row_heights if any(v is not None for v in row_heights) else None

            schemas.append(
                TableSchema(
                    index=g_idx,
                    rows=n_rows,
                    cols=n_cols,
                    preview=preview,
                    location=f"slide {s_idx}",
                    merges=merges,
                    column_widths_cm=col_widths_out,
                    row_heights_cm=row_heights_out,
                )
            )
        return schemas

    def get_cell(self, table_index: int, row: int, col: int) -> CellContent:
        table = self._get_table(table_index)
        n_rows, n_cols = self._dimensions(table)
        if row < 0 or col < 0 or row >= n_rows or col >= n_cols:
            raise CellOutOfBoundsError(
                f"cell ({row},{col}) out of bounds ({n_rows}x{n_cols})"
            )
        cell = table.cell(row, col)

        if cell.is_merge_origin:
            anchor = (row, col)
            span = (cell.span_height, cell.span_width)
            is_anchor = True
            source_cell = cell
        elif cell.is_spanned:
            anchor = self._find_merge_origin(table, row, col)
            source_cell = table.cell(*anchor)
            span = (source_cell.span_height, source_cell.span_width)
            is_anchor = False
        else:
            anchor = (row, col)
            span = (1, 1)
            is_anchor = True
            source_cell = cell

        paragraphs_text = [p.text for p in source_cell.text_frame.paragraphs]
        text = source_cell.text or ""

        # 셀 크기 힌트: anchor 위치부터 span 만큼의 column/row 합.
        a_r, a_c = anchor
        r_span, c_span = span
        try:
            cols_list = list(table.columns)
            rows_list = list(table.rows)
            width_emu = sum(
                getattr(cols_list[i], "width", 0) or 0
                for i in range(a_c, min(a_c + c_span, len(cols_list)))
            )
            height_emu = sum(
                getattr(rows_list[i], "height", 0) or 0
                for i in range(a_r, min(a_r + r_span, len(rows_list)))
            )
        except (IndexError, AttributeError):
            width_emu = 0
            height_emu = 0

        return CellContent(
            row=row,
            col=col,
            text=text,
            paragraphs=paragraphs_text,
            is_anchor=is_anchor,
            anchor=anchor,
            span=span,
            nested_table_indices=[],  # PPTX는 중첩 테이블 미지원
            width_cm=_emu_to_cm(width_emu),
            height_cm=_emu_to_cm(height_emu),
            char_count=len(text),
        )

    # ---- shape text (v0.8+) ----
    def get_shapes(
        self,
        slide_index: int | None = None,
        min_text_len: int = 1,
        max_preview: int = 40,
    ) -> list[ShapeInfo]:
        """표 외 shape (textbox / placeholder / 도형 텍스트) 수집.

        ``slide_index`` 는 1-based. None 이면 전체.
        ``min_text_len`` 미만의 텍스트는 제외 (0 을 주면 빈 shape 도 포함).
        """
        shapes_out: list[ShapeInfo] = []
        for s_idx, slide in enumerate(self._prs.slides, 1):
            if slide_index is not None and s_idx != slide_index:
                continue
            for shape in slide.shapes:
                if shape.has_table:
                    continue  # 표는 get_tables 로
                if not shape.has_text_frame:
                    continue
                text = (shape.text_frame.text or "").strip()
                if len(text) < min_text_len:
                    continue
                ph_type = None
                try:
                    ph = shape.placeholder_format
                    if ph is not None and ph.type is not None:
                        ph_type = str(ph.type).rsplit(".", 1)[-1]
                except ValueError:
                    pass
                kind = "placeholder" if ph_type else "text_box"
                shapes_out.append(
                    ShapeInfo(
                        slide_index=s_idx,
                        shape_id=shape.shape_id,
                        name=shape.name,
                        kind=kind,
                        has_text=bool(text),
                        text=text,
                        text_preview=text[:max_preview],
                        placeholder_type=ph_type,
                    )
                )
        return shapes_out

    def set_shape_text(
        self,
        slide_index: int,
        shape_id: int,
        text: str,
    ) -> str:
        """shape 의 텍스트를 text 로 교체. 기존 run-level 포맷 보존."""
        for s_idx, slide in enumerate(self._prs.slides, 1):
            if s_idx != slide_index:
                continue
            for shape in slide.shapes:
                if shape.shape_id != shape_id:
                    continue
                if not shape.has_text_frame:
                    raise ValueError(
                        f"shape {shape_id} on slide {slide_index} has no text frame"
                    )
                old = shape.text_frame.text or ""
                _set_text_frame_preserving_format(shape.text_frame, text)
                return old
        raise ValueError(
            f"shape not found: slide_index={slide_index}, shape_id={shape_id}"
        )

    # ---- charts (v0.17+) ----

    def _read_chart(
        self, chart: Any
    ) -> tuple[list[str], list[dict[str, Any]], bool, str | None]:
        """(categories, series, editable, warning) — 절대 예외를 밖으로 내지 않는다.

        get_charts 와 set_chart_data 의 before/after 스냅샷이 공용으로 사용.
        편집 불가 구조(콤보/다중레벨 카테고리/날짜축/scatter·bubble)는
        editable=False + 이유(warning)로 표시한다.
        """
        editable = True
        warning: str | None = None

        def _mark(msg: str) -> None:
            nonlocal editable, warning
            editable = False
            if warning is None:
                warning = msg

        cats: list[str] = []
        sers: list[dict[str, Any]] = []
        try:
            plots = list(chart.plots)
        except Exception as e:  # noqa: BLE001 — 미지원 차트도 목록엔 나와야 함
            return [], [], False, f"cannot read chart plots: {e}"
        if len(plots) != 1:
            _mark(f"combo chart with {len(plots)} plots — editing not supported")
        if plots:
            try:
                cat_obj = plots[0].categories
                if getattr(cat_obj, "depth", 1) > 1:
                    _mark("multi-level categories — editing not supported")
                raw = list(cat_obj)
                if any(isinstance(c, (datetime.date, datetime.datetime))
                       for c in raw):
                    _mark("date-axis categories — editing not supported")
                cats = ["" if c is None else str(c) for c in raw]
            except Exception as e:  # noqa: BLE001
                _mark(f"cannot read categories: {e}")
        try:
            ctname = chart.chart_type.name
            if ctname.startswith(("XY_", "BUBBLE")):
                _mark(f"{ctname} charts are read-only (v1)")
        except Exception:  # noqa: BLE001
            _mark("unknown chart type")
        try:
            for i, ser in enumerate(chart.series):
                try:
                    name = str(ser.name)
                except Exception:  # noqa: BLE001
                    name = f"Series {i + 1}"
                try:
                    values = list(ser.values)
                except Exception as e:  # noqa: BLE001
                    values = []
                    _mark(f"cannot read series values: {e}")
                sers.append({"name": name, "values": values})
        except Exception as e:  # noqa: BLE001
            _mark(f"cannot read series: {e}")
        return cats, sers, editable, warning

    def get_charts(self, slide_index: int | None = None) -> list[ChartInfo]:
        """차트 목록 + 데이터. ``slide_index`` 는 1-based, None 이면 전체.

        차트는 get_tables/get_shapes 에 나타나지 않으므로 차트 작업의 진입점은
        항상 이 메서드다. 그룹 shape 안의 차트도 수집한다.
        """
        out: list[ChartInfo] = []
        for s_idx, slide in enumerate(self._prs.slides, 1):
            if slide_index is not None and s_idx != slide_index:
                continue
            for shape in _iter_shapes_recursive(slide.shapes):
                if not getattr(shape, "has_chart", False):
                    continue
                chart = shape.chart
                try:
                    ctype = chart.chart_type.name
                except Exception:  # noqa: BLE001
                    ctype = "UNKNOWN"
                title = None
                try:
                    if chart.has_title:
                        title = chart.chart_title.text_frame.text or None
                except Exception:  # noqa: BLE001
                    pass
                cats, sers, editable, warning = self._read_chart(chart)
                out.append(ChartInfo(
                    slide_index=s_idx,
                    shape_id=shape.shape_id,
                    name=shape.name,
                    chart_type=ctype,
                    title=title,
                    categories=cats,
                    series=sers,
                    editable=editable,
                    warning=warning,
                ))
        return out

    def _find_chart(self, slide_index: int, shape_id: int) -> Any:
        """(slide_index, shape_id) 로 chart 객체 탐색. 실패는 ValueError."""
        slides = list(self._prs.slides)
        if slide_index < 1 or slide_index > len(slides):
            raise ValueError(
                f"slide_index {slide_index} out of range 1..{len(slides)}"
            )
        for shape in _iter_shapes_recursive(slides[slide_index - 1].shapes):
            if shape.shape_id != shape_id:
                continue
            if not getattr(shape, "has_chart", False):
                raise ValueError(
                    f"shape {shape_id} on slide {slide_index} is not a chart "
                    f"(use get_charts to list chart shape_ids)"
                )
            return shape.chart
        raise ValueError(
            f"chart not found: slide_index={slide_index}, shape_id={shape_id}"
        )

    def set_chart_data(
        self,
        slide_index: int,
        shape_id: int,
        *,
        categories: list[str] | None = None,
        series: list[dict[str, Any]] | None = None,
        set_points: list[dict[str, Any]] | None = None,
        title: str | None = None,
    ) -> dict[str, Any]:
        """차트 데이터 편집 — read → 수정 → replace_data (서식 보존).

        모드/검증 규칙은 base.DocumentAdapter.set_chart_data docstring 참조.
        """
        chart = self._find_chart(slide_index, shape_id)
        cats, sers, editable, warning = self._read_chart(chart)

        if series is None and set_points is None and title is None:
            raise ValueError(
                "nothing to change: pass series, set_points, or title. "
                "변경할 내용이 없습니다 — series/set_points/title 중 하나를 전달하세요."
            )
        if series is not None and set_points is not None:
            raise ValueError(
                "pass either series or set_points, not both. "
                "series 와 set_points 는 동시에 줄 수 없습니다."
            )
        if set_points is not None and categories is not None:
            raise ValueError(
                "categories cannot be combined with set_points — use the "
                "full-replace mode (categories + series) to change categories. "
                "카테고리 변경은 전체 교체 모드(categories+series)를 사용하세요."
            )
        if (series is not None or set_points is not None) and not editable:
            raise NotImplementedForFormat(
                f"this chart cannot be edited: {warning}. "
                f"이 차트는 데이터 편집을 지원하지 않습니다: {warning}"
            )

        before = {
            "categories": list(cats),
            "series": [dict(s, values=list(s["values"])) for s in sers],
        }

        new_cats: list[str] | None = None
        new_sers: list[dict[str, Any]] | None = None
        if series is not None:
            new_cats = ([str(c) for c in categories]
                        if categories is not None else list(cats))
            if not new_cats:
                raise ValueError(
                    "chart has no categories — pass categories explicitly. "
                    "카테고리가 없습니다 — categories 를 명시하세요."
                )
            new_sers = _validate_series(series, new_cats)
        elif set_points is not None:
            if not isinstance(set_points, (list, tuple)) or not set_points:
                raise ValueError(
                    "set_points must be a non-empty list of "
                    "{series, category, value}."
                )
            new_cats = list(cats)
            new_sers = [dict(s, values=list(s["values"])) for s in sers]
            ser_names = [s["name"] for s in new_sers]
            for j, pt in enumerate(set_points):
                if not isinstance(pt, dict):
                    raise ValueError(f"set_points[{j}] must be a dict")
                si = _resolve_index(pt.get("series"), ser_names, "series")
                ci = _resolve_index(pt.get("category"), new_cats, "category")
                new_sers[si]["values"][ci] = _to_number(pt.get("value"))

        if new_sers is not None:
            cd = CategoryChartData()
            cd.categories = new_cats
            for s in new_sers:
                cd.add_series(s["name"], tuple(s["values"]))
            chart.replace_data(cd)

        if title is not None:
            chart.has_title = True
            chart.chart_title.text_frame.text = title

        a_cats, a_sers, _, _ = self._read_chart(chart)
        try:
            ctype = chart.chart_type.name
        except Exception:  # noqa: BLE001
            ctype = "UNKNOWN"
        result: dict[str, Any] = {
            "slide_index": slide_index,
            "shape_id": shape_id,
            "chart_type": ctype,
            "before": before,
            "after": {"categories": a_cats, "series": a_sers},
        }
        if title is not None:
            result["title"] = title
        return result

    def add_chart(
        self,
        slide_index: int,
        chart_type: str,
        *,
        categories: list[str],
        series: list[dict[str, Any]],
        title: str | None = None,
        x_cm: float | None = None,
        y_cm: float | None = None,
        width_cm: float | None = None,
        height_cm: float | None = None,
    ) -> dict[str, Any]:
        """슬라이드에 새 차트 추가. 위치 생략 시 제목 아래 전폭 배치."""
        slides = list(self._prs.slides)
        if slide_index < 1 or slide_index > len(slides):
            raise ValueError(
                f"slide_index {slide_index} out of range 1..{len(slides)}"
            )
        slide = slides[slide_index - 1]

        key = str(chart_type).strip().lower()
        xl_type = _CHART_TYPE_MAP.get(key)
        if xl_type is None:
            raise ValueError(
                f"unsupported chart_type {chart_type!r}. "
                f"supported: {sorted(_CHART_TYPE_MAP)}. "
                f"지원하지 않는 차트 타입입니다 — 지원 목록에서 고르세요."
            )
        if not categories:
            raise ValueError(
                "categories must be a non-empty list. "
                "categories 는 비어있지 않은 리스트여야 합니다."
            )
        cats = [str(c) for c in categories]
        sers = _validate_series(series, cats)

        x = x_cm if x_cm is not None else _CHART_DEFAULT_X_CM
        y = y_cm if y_cm is not None else _CHART_DEFAULT_Y_CM
        slide_w_cm = int(self._prs.slide_width or 0) / _EMU_PER_CM
        slide_h_cm = int(self._prs.slide_height or 0) / _EMU_PER_CM
        w = width_cm if width_cm is not None else max(slide_w_cm - 2 * x, 2.0)
        h = height_cm if height_cm is not None else max(slide_h_cm - y - 1.0, 2.0)

        cd = CategoryChartData()
        cd.categories = cats
        for s in sers:
            cd.add_series(s["name"], tuple(s["values"]))
        gf = slide.shapes.add_chart(
            xl_type,
            Emu(int(x * _EMU_PER_CM)), Emu(int(y * _EMU_PER_CM)),
            Emu(int(w * _EMU_PER_CM)), Emu(int(h * _EMU_PER_CM)),
            cd,
        )
        chart = gf.chart
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        if len(sers) > 1:
            try:
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
            except Exception:  # noqa: BLE001 — 범례는 부가 기능, 실패해도 차트는 유효
                pass
        return {
            "slide_index": slide_index,
            "shape_id": gf.shape_id,
            "chart_type": xl_type.name,
            "categories": cats,
            "series_names": [s["name"] for s in sers],
        }

    # ---- slides (v0.17+) ----

    def get_slides(self) -> list[SlideInfo]:
        """슬라이드 개요 목록 — 복제/편집할 '양식 페이지' 를 고르는 눈."""
        out: list[SlideInfo] = []
        for s_idx, slide in enumerate(self._prs.slides, 1):
            shapes = list(_iter_shapes_recursive(slide.shapes))
            n_tables = sum(
                1 for s in shapes if getattr(s, "has_table", False))
            n_charts = sum(
                1 for s in shapes if getattr(s, "has_chart", False))
            text_shapes = [
                s for s in shapes
                if getattr(s, "has_text_frame", False)
                and (s.text_frame.text or "").strip()
            ]
            title: str | None = None
            for s in shapes:
                if not getattr(s, "is_placeholder", False):
                    continue
                try:
                    if (s.placeholder_format.idx == 0
                            and getattr(s, "has_text_frame", False)):
                        title = (s.text_frame.text or "").strip()[:40] or None
                        break
                except ValueError:
                    continue
            if title is None and text_shapes:
                title = (text_shapes[0].text_frame.text or "").strip()[:40] or None
            try:
                layout_name = slide.slide_layout.name or ""
            except Exception:  # noqa: BLE001
                layout_name = ""
            out.append(SlideInfo(
                slide_index=s_idx,
                layout_name=layout_name,
                title=title,
                shape_count=len(shapes),
                table_count=n_tables,
                chart_count=n_charts,
                text_shape_count=len(text_shapes),
                texts_preview=[
                    (s.text_frame.text or "").strip()[:40]
                    for s in text_shapes[:5]
                ],
            ))
        return out

    def duplicate_slide(
        self,
        source_slide_index: int,
        at: int | None = None,
    ) -> dict[str, Any]:
        """양식 슬라이드 복제 — 같은 레이아웃의 빈 슬라이드에 shape XML 을
        deepcopy 하고 관계(rId)를 재매핑한다. 차트 part 는 내장 워크북까지
        독립 복제해 복제본 편집이 원본을 오염시키지 않는다.

        노트 슬라이드는 복제하지 않는다 (v1 계약).
        """
        slides = list(self._prs.slides)
        n = len(slides)
        if source_slide_index < 1 or source_slide_index > n:
            raise ValueError(
                f"source_slide_index {source_slide_index} out of range 1..{n}"
            )
        if at is not None and (at < 1 or at > n + 1):
            raise ValueError(f"at {at} out of range 1..{n + 1}")

        source = slides[source_slide_index - 1]
        new_slide = self._prs.slides.add_slide(source.slide_layout)

        # add_slide 가 레이아웃에서 자동 생성한 placeholder 제거 — 원본 shape 로 대체
        for shp in list(new_slide.shapes):
            shp._element.getparent().remove(shp._element)

        # 원본 slide 의 관계 재생성 + {구 rId → 새 rId} 맵
        id_map: dict[str, str] = {}
        for rid in sorted(source.part.rels, key=lambda r: int(r[3:])):
            rel = source.part.rels[rid]
            if rel.reltype in (RT.SLIDE_LAYOUT, RT.NOTES_SLIDE):
                continue  # layout 은 add_slide 가 연결함. 노트는 v1 미복제.
            if rel.is_external:
                id_map[rid] = new_slide.part.rels.get_or_add_ext_rel(
                    rel.reltype, rel.target_ref
                )
            else:
                target = rel.target_part
                if rel.reltype == RT.CHART:
                    target = _clone_part(target)
                id_map[rid] = new_slide.part.relate_to(target, rel.reltype)

        # shape XML deepcopy + r:* 참조 재매핑
        r_prefix = f"{{{_R_NS}}}"
        for shape in source.shapes:
            el = deepcopy(shape._element)
            for node in el.iter():
                for attr, val in list(node.attrib.items()):
                    if attr.startswith(r_prefix) and val in id_map:
                        node.set(attr, id_map[val])
            new_slide.shapes._spTree.append(el)

        # 위치 지정: sldIdLst 에서 마지막(새) sldId 를 원하는 위치로 이동
        new_index = n + 1
        if at is not None and at <= n:
            sldIdLst = self._prs.slides._sldIdLst
            ids = list(sldIdLst)
            new_id = ids[-1]
            sldIdLst.remove(new_id)
            sldIdLst.insert(at - 1, new_id)
            new_index = at

        # 좌표 피드백 — 삽입 *후* 좌표계로 재계산 (create_document 관례).
        # preview(행 라벨 포함)를 함께 담는다 — LLM 이 복제본을 채울 때 행 매핑을
        # 추측하다 off-by-one 으로 어긋나는 실패 패턴 방지 (라벨을 보고 좌표 결정).
        previews = {
            t.index: t.preview
            for t in self.get_tables(preview_rows=10, max_cell_len=30)
        }
        tables = []
        for g_idx, s_idx, table in self._iter_tables():
            if s_idx != new_index:
                continue
            rows, cols = self._dimensions(table)
            tables.append({
                "table_index": g_idx,
                "rows": rows,
                "cols": cols,
                "preview": previews.get(g_idx),
            })
        charts = [
            {"shape_id": c.shape_id, "chart_type": c.chart_type}
            for c in self.get_charts(slide_index=new_index)
        ]
        text_shapes = [
            {"shape_id": s.shape_id, "name": s.name,
             "text_preview": s.text_preview}
            for s in self.get_shapes(slide_index=new_index, min_text_len=0)
        ]
        result: dict[str, Any] = {
            "source_slide_index": source_slide_index,
            "new_slide_index": new_index,
            "slide_count": n + 1,
            "tables": tables,
            "charts": charts,
            "text_shapes": text_shapes,
        }
        if at is not None and at <= n:
            result["warning"] = (
                "슬라이드가 중간에 삽입되어 뒤쪽 표들의 table_index 가 "
                "변경되었습니다. 이전 inspect 결과의 table_index 를 신뢰하지 "
                "말고 이 반환값과 재-inspect 결과를 사용하세요."
            )
        return result

    # ---- shape copy (v0.18+) ----

    def _find_table_shape(self, table_index: int) -> tuple[int, Any]:
        """전역 flat table_index → (slide_index_1based, graphicFrame shape)."""
        g_idx = 0
        for s_idx, slide in enumerate(self._prs.slides, 1):
            for shape in slide.shapes:
                if getattr(shape, "has_table", False):
                    if g_idx == table_index:
                        return s_idx, shape
                    g_idx += 1
        raise TableIndexError(f"PPTX table index {table_index} not found")

    def copy_shape(
        self,
        target_slide_index: int,
        *,
        source_slide_index: int | None = None,
        shape_id: int | None = None,
        table_index: int | None = None,
        x_cm: float | None = None,
        y_cm: float | None = None,
        clear_values: bool = False,
    ) -> dict[str, Any]:
        """표/차트/텍스트박스 shape 하나를 다른 슬라이드로 복사 (서식 유지).

        레시피는 duplicate_slide 의 shape 단위 축소판:
        deepcopy → (placeholder 면 위치·크기 실측 고정 후 <p:ph> 제거) →
        참조된 관계(rId)만 대상 슬라이드에 재생성 (차트는 part 독립 복제) →
        cNvPr id 를 대상 슬라이드 내 유일값으로 재부여 → spTree append →
        위치 적용 → (옵션) 값 비우기.
        """
        slides = list(self._prs.slides)
        n = len(slides)
        if target_slide_index < 1 or target_slide_index > n:
            raise ValueError(
                f"target_slide_index {target_slide_index} out of range 1..{n}"
            )

        # ---- 원본 shape 해소 ----
        if (shape_id is None) == (table_index is None):
            raise ValueError(
                "pass exactly one of table_index (for tables) or "
                "source_slide_index+shape_id (for charts/text shapes). "
                "표는 table_index 로, 차트/텍스트박스는 "
                "source_slide_index+shape_id 로 지정하세요 (동시 지정 불가)."
            )
        if table_index is not None:
            src_s_idx, src_shape = self._find_table_shape(table_index)
            if source_slide_index is not None and source_slide_index != src_s_idx:
                raise ValueError(
                    f"table {table_index} is on slide {src_s_idx}, "
                    f"not slide {source_slide_index}"
                )
        else:
            if source_slide_index is None:
                raise ValueError(
                    "shape_id requires source_slide_index. "
                    "shape_id 지정 시 source_slide_index 도 필요합니다."
                )
            if source_slide_index < 1 or source_slide_index > n:
                raise ValueError(
                    f"source_slide_index {source_slide_index} out of range 1..{n}"
                )
            src_s_idx = source_slide_index
            src_slide_obj = slides[src_s_idx - 1]
            src_shape = None
            for sh in src_slide_obj.shapes:
                if sh.shape_id == shape_id:
                    src_shape = sh
                    break
            if src_shape is None:
                for sh in _iter_shapes_recursive(src_slide_obj.shapes):
                    if sh.shape_id == shape_id:
                        raise ValueError(
                            f"shape {shape_id} is inside a group shape — "
                            f"copy the group itself instead. "
                            f"그룹 내부 shape 는 그룹 전체의 shape_id 로 복사하세요."
                        )
                raise ValueError(
                    f"shape not found: slide_index={src_s_idx}, shape_id={shape_id}"
                )

        src_slide = slides[src_s_idx - 1]
        target_slide = slides[target_slide_index - 1]

        kind = (
            "table" if getattr(src_shape, "has_table", False)
            else "chart" if getattr(src_shape, "has_chart", False)
            else "text" if getattr(src_shape, "has_text_frame", False)
            else "other"
        )

        # placeholder 는 위치/크기를 레이아웃에서 상속할 수 있으므로, 복사 전에
        # 실측값(python-pptx 가 상속 해소한 값)을 확보해 둔다.
        eff: dict[str, Any] = {}
        for attr in ("left", "top", "width", "height"):
            try:
                eff[attr] = getattr(src_shape, attr)
            except Exception:  # noqa: BLE001
                eff[attr] = None

        el = deepcopy(src_shape._element)

        # <p:ph> 제거 — placeholder 복사본이 대상 슬라이드의 placeholder 와
        # idx 충돌하는 것을 방지 (일반 shape 로 전환, 위치는 아래서 고정).
        ph_els = list(el.iter(f"{{{_P_NS}}}ph"))
        for ph_el in ph_els:
            parent = ph_el.getparent()
            if parent is not None:
                parent.remove(ph_el)

        # 참조된 관계(rId)만 대상 슬라이드에 재생성 + 재매핑
        r_prefix = f"{{{_R_NS}}}"
        referenced: set[str] = set()
        for node in el.iter():
            for attr, val in node.attrib.items():
                if attr.startswith(r_prefix):
                    referenced.add(val)
        id_map: dict[str, str] = {}
        for rid in sorted(referenced):
            if rid not in src_slide.part.rels:
                continue
            rel = src_slide.part.rels[rid]
            if rel.is_external:
                id_map[rid] = target_slide.part.rels.get_or_add_ext_rel(
                    rel.reltype, rel.target_ref
                )
            else:
                target_part = rel.target_part
                if rel.reltype == RT.CHART:
                    target_part = _clone_part(target_part)
                id_map[rid] = target_slide.part.relate_to(
                    target_part, rel.reltype
                )
        for node in el.iter():
            for attr, val in list(node.attrib.items()):
                if attr.startswith(r_prefix) and val in id_map:
                    node.set(attr, id_map[val])

        # cNvPr id 재부여 — 대상 슬라이드 내 유일 보장 (그룹이면 내부까지)
        used_ids = {
            int(c.get("id"))
            for c in target_slide._element.iter(f"{{{_P_NS}}}cNvPr")
            if (c.get("id") or "").isdigit()
        }
        next_id = max(used_ids, default=1) + 1
        new_ids: list[int] = []
        for c in el.iter(f"{{{_P_NS}}}cNvPr"):
            c.set("id", str(next_id))
            new_ids.append(next_id)
            next_id += 1
        if not new_ids:
            raise NotImplementedForFormat(
                "copied element has no cNvPr id — unsupported shape structure"
            )
        new_shape_id = new_ids[0]

        target_slide.shapes._spTree.append(el)

        new_shape = None
        for sh in target_slide.shapes:
            if sh.shape_id == new_shape_id:
                new_shape = sh
                break

        # 위치/크기 적용 — x/y 지정 시 그 위치, 아니면 원본 실측 위치.
        # (placeholder 였던 shape 는 xfrm 이 없을 수 있어 실측값으로 고정.)
        if new_shape is not None:
            left = (Emu(int(x_cm * _EMU_PER_CM)) if x_cm is not None
                    else eff["left"])
            top = (Emu(int(y_cm * _EMU_PER_CM)) if y_cm is not None
                   else eff["top"])
            try:
                if left is not None:
                    new_shape.left = left
                if top is not None:
                    new_shape.top = top
                if getattr(new_shape, "width", None) is None \
                        and eff["width"] is not None:
                    new_shape.width = eff["width"]
                if getattr(new_shape, "height", None) is None \
                        and eff["height"] is not None:
                    new_shape.height = eff["height"]
            except (AttributeError, ValueError):
                pass  # 위치 미지원 shape — 복사 자체는 유효

        # 값 비우기 (서식/구조 유지) — 표 셀 run, 텍스트 run. 차트는 무의미.
        cleared = False
        if clear_values and new_shape is not None and kind in ("table", "text"):
            a = f"{{{_A_NS}}}"
            if kind == "table":
                for tc in new_shape._element.iter(f"{a}tc"):
                    self._blank_copied_pptx_tc(tc, a)
            else:
                for t_el in new_shape._element.iter(f"{a}t"):
                    t_el.text = ""
            cleared = True

        # ---- 좌표 피드백 ----
        result: dict[str, Any] = {
            "source_slide_index": src_s_idx,
            "target_slide_index": target_slide_index,
            "shape_id": new_shape_id,
            "kind": kind,
            "values_cleared": cleared,
        }
        if kind == "table":
            g = 0
            new_tidx = None
            total_tables = 0
            for s_i, slide in enumerate(self._prs.slides, 1):
                for sh in slide.shapes:
                    if getattr(sh, "has_table", False):
                        if s_i == target_slide_index \
                                and sh.shape_id == new_shape_id:
                            new_tidx = g
                        g += 1
            total_tables = g
            previews = {
                t.index: t.preview
                for t in self.get_tables(preview_rows=10, max_cell_len=30)
            }
            result["table_index"] = new_tidx
            if new_tidx is not None:
                result["preview"] = previews.get(new_tidx)
            if new_tidx is not None and new_tidx < total_tables - 1:
                result["warning"] = (
                    "새 표가 중간 순번에 추가되어 뒤쪽 표들의 전역 table_index "
                    "가 변경되었습니다. 이 반환값의 table_index 를 사용하세요."
                )
        elif kind == "chart":
            for c in self.get_charts(slide_index=target_slide_index):
                if c.shape_id == new_shape_id:
                    result["chart_type"] = c.chart_type
                    break
            result["hint"] = (
                "복사본 수치 변경은 set_chart_data("
                f"slide_index={target_slide_index}, shape_id={new_shape_id}) 로."
            )
        return result

    # ---- editing ----
    def render_template(self, context: dict[str, Any], *,
                        on_missing: str = "blank") -> dict[str, list[str]]:
        """paragraph 단위로 {{key}}를 치환. run이 쪼개진 경우를 처리하기 위해
        paragraph 전체 텍스트를 재조립 후 첫 run에 담는다 (서식 일부 손실 가능).
        누락 키 처리는 on_missing 정책을 따른다 (base 참조)."""
        report = self._render_report(self.get_placeholders(), context, on_missing)
        for tf in self._iter_text_frames():
            for para in tf.paragraphs:
                full_text = "".join(run.text for run in para.runs)
                if not _has_template(full_text):
                    continue
                rendered = self._render_text_block(full_text, context, on_missing)
                if para.runs:
                    para.runs[0].text = rendered
                    for run in para.runs[1:]:
                        run.text = ""
        return report

    def _get_table(self, table_index: int):
        for g_idx, _, table in self._iter_tables():
            if g_idx == table_index:
                return table
        raise TableIndexError(f"PPTX table index {table_index} not found")

    def set_cell(
        self,
        table_index: int,
        row: int,
        col: int,
        value: str,
        *,
        allow_merge_redirect: bool = False,
    ) -> str:
        table = self._get_table(table_index)
        cell, _, _, _ = self._resolve_anchor_cell(
            table, row, col, allow_merge_redirect=allow_merge_redirect
        )
        old = cell.text
        _set_text_frame_preserving_format(cell.text_frame, value)
        return old

    def append_to_cell(
        self,
        table_index: int,
        row: int,
        col: int,
        value: str,
        separator: str = "  ",
        *,
        allow_merge_redirect: bool = False,
    ) -> str:
        table = self._get_table(table_index)
        cell, _, _, _ = self._resolve_anchor_cell(
            table, row, col, allow_merge_redirect=allow_merge_redirect
        )
        old = cell.text
        new_value = f"{old}{separator}{value}" if old else value
        _set_text_frame_preserving_format(cell.text_frame, new_value)
        return old

    def append_row(self, table_index: int, values: list[str]) -> None:
        """표 끝에 새 행 추가 — 마지막 <a:tr> 을 deepcopy 후 텍스트만 비움.

        python-pptx 공식 add_row 는 없지만 OOXML 스키마상 <a:tr> 을 붙이는 것만으로
        행이 추가된다. 마지막 행의 셀 구조 (gridSpan/rowSpan/hMerge/vMerge, tcPr
        스타일) 를 그대로 상속해 이전 행과 동일한 서식의 빈 행이 생긴다.

        제약:
          - 마지막 행이 위 행의 rowSpan 영역에 속하면 (vMerge="1" 또는 rowSpan>1 셀
            존재) 복제 시 교차 병합이 오동작하므로 ``NotImplementedForFormat``.
        """
        table = self._get_table(table_index)
        tbl_elem = table._tbl  # lxml <a:tbl>

        a = f"{{{_A_NS}}}"
        trs = tbl_elem.findall(f"{a}tr")
        if not trs:
            raise NotImplementedForFormat("cannot append row to empty PPTX table")

        last_row = trs[-1]
        for tc in last_row.findall(f"{a}tc"):
            if tc.get("vMerge") == "1":
                raise NotImplementedForFormat(
                    "last row participates in a cross-row merge (vMerge); "
                    "append_row is not safe for this table."
                )
            try:
                rs = int(tc.get("rowSpan", "1"))
            except (TypeError, ValueError):
                rs = 1
            if rs > 1:
                raise NotImplementedForFormat(
                    "last row contains a rowSpan anchor that extends past the table; "
                    "append_row is not safe for this table."
                )

        new_row = deepcopy(last_row)
        # 기존 run/paragraph 구조는 유지하고 <a:t>.text 만 비움 (스타일 보존)
        for tc in new_row.findall(f"{a}tc"):
            txBody = tc.find(f"{a}txBody")
            if txBody is None:
                continue
            for p in txBody.findall(f"{a}p"):
                for r_el in p.findall(f"{a}r"):
                    for t_el in r_el.findall(f"{a}t"):
                        t_el.text = ""
        tbl_elem.append(new_row)

        new_row_idx = len(trs)  # 새 행 인덱스 (append 전 길이 = 새 행 position)
        n_cols = len(list(table.columns))
        for i, value in enumerate(values):
            if i >= n_cols:
                break
            try:
                self.set_cell(table_index, new_row_idx, i, value)
            except MergedCellWriteError:
                # 복제로 상속된 병합의 non-anchor 좌표는 자연히 스킵
                continue

    # ---- 위치 지정 행/열 삽입 (v0.14+) ----

    @staticmethod
    def _blank_copied_pptx_tc(tc, a: str) -> None:
        """deepcopy 된 <a:tc> 의 run 텍스트만 비움 (스타일/구조 유지)."""
        txBody = tc.find(f"{a}txBody")
        if txBody is None:
            return
        for p in txBody.findall(f"{a}p"):
            for r_el in p.findall(f"{a}r"):
                for t_el in r_el.findall(f"{a}t"):
                    t_el.text = ""

    def insert_row(
        self,
        table_index: int,
        values: list[str],
        at_row: int | None = None,
    ) -> int:
        """지정 위치에 새 행 삽입 — 인접 <a:tr> deepcopy 로 서식 상속.

        DrawingML 표는 모든 grid 위치에 물리 <a:tc> 가 존재하므로 (병합
        continuation 은 hMerge/vMerge="1"), 템플릿 행 복사 + rowSpan/vMerge
        리셋으로 안전하게 독립 행을 만든다.
        """
        table = self._get_table(table_index)
        tbl_elem = table._tbl
        a = f"{{{_A_NS}}}"
        trs = tbl_elem.findall(f"{a}tr")
        n_rows = len(trs)
        if n_rows == 0:
            raise NotImplementedForFormat(
                "cannot insert a row into an empty PPTX table"
            )
        if at_row is None:
            at_row = n_rows
        if at_row < 0 or at_row > n_rows:
            raise CellOutOfBoundsError(f"at_row {at_row} out of range 0..{n_rows}")

        # 세로 병합 경계 가드: 삽입 위치 행에 vMerge continuation 이 있으면
        # 위 행에서 내려오는 병합이 경계를 가로지른다는 뜻.
        if 0 < at_row < n_rows:
            for tc in trs[at_row].findall(f"{a}tc"):
                if tc.get("vMerge") == "1":
                    raise NotImplementedForFormat(
                        f"insertion point row {at_row} crosses a vertical "
                        f"merge; inserting here would split the merged region."
                    )
        if at_row == n_rows:
            # 맨 끝 삽입: 기존 append_row 와 동일한 마지막 행 안전성 검사
            for tc in trs[-1].findall(f"{a}tc"):
                if tc.get("vMerge") == "1":
                    raise NotImplementedForFormat(
                        "last row participates in a cross-row merge (vMerge); "
                        "inserting at the end is not safe for this table."
                    )

        template_idx = at_row if at_row < n_rows else n_rows - 1
        new_row = deepcopy(trs[template_idx])
        for tc in new_row.findall(f"{a}tc"):
            # 새 행은 세로 병합에 참여하지 않는 독립 행
            if tc.get("rowSpan"):
                del tc.attrib["rowSpan"]
            if tc.get("vMerge"):
                del tc.attrib["vMerge"]
            self._blank_copied_pptx_tc(tc, a)

        if at_row < n_rows:
            trs[at_row].addprevious(new_row)
        else:
            trs[-1].addnext(new_row)

        n_cols = len(list(table.columns))
        for i, value in enumerate(values):
            if i >= n_cols:
                break
            if not value:
                continue
            try:
                self.set_cell(table_index, at_row, i, value)
            except MergedCellWriteError:
                continue
        return at_row

    def insert_column(
        self,
        table_index: int,
        values: list[str],
        at_col: int | None = None,
    ) -> int:
        """지정 위치에 새 열 삽입 — 행별 인접 <a:tc> deepcopy 로 서식 상속.

        <a:tblGrid> 에 gridCol 을 삽입하고 각 행에 새 <a:tc> 를 끼워 넣는다.
        표 전체 폭 유지를 위해 gridCol 폭을 비례 축소해 새 열 폭을 흡수한다.
        """
        table = self._get_table(table_index)
        tbl_elem = table._tbl
        a = f"{{{_A_NS}}}"
        trs = tbl_elem.findall(f"{a}tr")
        grid_el = tbl_elem.find(f"{a}tblGrid")
        grid_cols = grid_el.findall(f"{a}gridCol") if grid_el is not None else []
        n_rows, n_cols = len(trs), len(grid_cols)
        if n_rows == 0 or n_cols == 0:
            raise NotImplementedForFormat(
                "cannot insert a column into an empty PPTX table"
            )
        if at_col is None:
            at_col = n_cols
        if at_col < 0 or at_col > n_cols:
            raise CellOutOfBoundsError(f"at_col {at_col} out of range 0..{n_cols}")

        # 가로 병합 경계 가드: 삽입 위치 열에 hMerge continuation 이 있으면 거부
        if 0 < at_col < n_cols:
            for tr in trs:
                tcs = tr.findall(f"{a}tc")
                if at_col < len(tcs) and tcs[at_col].get("hMerge") == "1":
                    raise NotImplementedForFormat(
                        f"insertion point column {at_col} crosses a horizontal "
                        f"merge; inserting here would split the merged region."
                    )

        template_col = at_col - 1 if at_col > 0 else 0

        # gridCol 삽입 + 폭 비례 재배분 (EMU)
        widths: list[int | None] = []
        for gc in grid_cols:
            try:
                widths.append(int(gc.get("w", "")))
            except (TypeError, ValueError):
                widths.append(None)
        new_gc = deepcopy(grid_cols[template_col])
        scale = 1.0
        if all(w is not None and w > 0 for w in widths):
            total = sum(widths)  # type: ignore[arg-type]
            new_w = widths[template_col]
            scale = total / (total + new_w)  # type: ignore[operator]
            for gc, w in zip(grid_cols, widths):
                gc.set("w", str(int(round(w * scale))))  # type: ignore[operator]
            new_gc.set("w", str(int(round(new_w * scale))))  # type: ignore[operator]
        if at_col < n_cols:
            grid_cols[at_col].addprevious(new_gc)
        else:
            grid_cols[-1].addnext(new_gc)

        # 행별 새 셀 삽입
        for tr in trs:
            tcs = tr.findall(f"{a}tc")
            if not tcs:
                continue
            template_tc = tcs[min(template_col, len(tcs) - 1)]
            new_tc = deepcopy(template_tc)
            for attr in ("gridSpan", "hMerge", "rowSpan", "vMerge"):
                if new_tc.get(attr):
                    del new_tc.attrib[attr]
            self._blank_copied_pptx_tc(new_tc, a)
            if at_col < len(tcs):
                tcs[at_col].addprevious(new_tc)
            else:
                tcs[-1].addnext(new_tc)

        for r, value in enumerate(values):
            if r >= n_rows:
                break
            if not value:
                continue
            try:
                self.set_cell(table_index, r, at_col, value)
            except (MergedCellWriteError, CellOutOfBoundsError):
                continue
        return at_col


def _set_text_frame_preserving_format(text_frame, value: str) -> None:
    """Write ``value`` into ``text_frame`` without losing run-level formatting.

    ``python-pptx`` exposes ``cell.text = value`` (which proxies to the text
    frame) but the setter deletes every run and replaces them with a single
    default-styled run. This destroys two kinds of formatting:

    1. **Runs that already exist** — font family, size, bold, color, etc.
    2. **Empty paragraphs that hold an ``<a:endParaRPr>``**, which is where
       PowerPoint stores the "what would the next character look like" run
       properties for an otherwise empty cell. Real-world templates put font
       information here so that the cell looks right even before any text
       is typed.

    Strategy:

    - If the paragraph already has runs, reuse the first one and blank the
      rest (simple case that covers pre-filled cells).
    - Otherwise, build a new ``<a:r>`` manually and clone ``<a:endParaRPr>``
      into its ``<a:rPr>`` so the empty-cell font survives.

    Paragraph comparison uses index, not identity, because python-pptx
    returns a fresh Python wrapper on every ``paragraphs`` access, which
    would cause a naive ``para is first_para`` check to always be False
    and blank the run we just populated.
    """
    paragraphs = list(text_frame.paragraphs)
    first_idx = next((i for i, para in enumerate(paragraphs) if para.runs), None)

    if first_idx is not None:
        first_para = paragraphs[first_idx]
        first_para.runs[0].text = value
        for run in first_para.runs[1:]:
            run.text = ""
        for i, para in enumerate(paragraphs):
            if i == first_idx:
                continue
            for run in para.runs:
                run.text = ""
        return

    target_para = paragraphs[0] if paragraphs else None
    if target_para is None:
        text_frame.text = value
        return

    p_el = target_para._p
    end_rpr = p_el.find(f"{{{_A_NS}}}endParaRPr")

    r_el = etree.SubElement(p_el, f"{{{_A_NS}}}r")
    if end_rpr is not None:
        rpr = deepcopy(end_rpr)
        rpr.tag = f"{{{_A_NS}}}rPr"
        r_el.insert(0, rpr)
    t_el = etree.SubElement(r_el, f"{{{_A_NS}}}t")
    t_el.text = value

    for i, para in enumerate(paragraphs):
        if i == 0:
            continue
        for run in para.runs:
            run.text = ""
