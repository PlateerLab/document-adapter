"""로컬 PPTX 파일을 스캔해서 fixture로 적합한 것 선별.

각 파일마다 뽑는 정보:
  - 파일 크기
  - 슬라이드 수
  - 표 개수
  - 병합 셀 개수 (anchor 기준)
  - 중첩 구조 유무 (shape 내부 shape 등)

우선순위가 높은 것 (다양성 + 표 많음 + 병합 셀 있음) 상위 N개를 후보로 출력.
"""
from __future__ import annotations

import subprocess
import sys
from pathlib import Path

SEARCH_ROOTS = [
    Path.home() / "Documents",
    Path.home() / "Desktop",
    Path.home() / "Downloads",
]


def find_pptx() -> list[Path]:
    results: list[Path] = []
    for root in SEARCH_ROOTS:
        if not root.exists():
            continue
        for p in root.rglob("*.pptx"):
            # 임시 파일 제외 (~$로 시작)
            if p.name.startswith("~$"):
                continue
            # node_modules 등 제외
            if "node_modules" in p.parts:
                continue
            results.append(p)
    return results


def analyze(path: Path) -> dict:
    """python-pptx로 열어 구조 메트릭 추출. 실패 시 error 필드."""
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx not installed"}

    try:
        prs = Presentation(str(path))
        slide_count = len(prs.slides)
        table_count = 0
        merge_count = 0
        max_rows = 0
        max_cols = 0
        total_cells = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table_count += 1
                    tbl = shape.table
                    rows = len(tbl.rows)
                    cols = len(tbl.columns)
                    max_rows = max(max_rows, rows)
                    max_cols = max(max_cols, cols)
                    for row in tbl.rows:
                        for cell in row.cells:
                            total_cells += 1
                            if getattr(cell, "is_merge_origin", False):
                                sh = getattr(cell, "span_height", 1) or 1
                                sw = getattr(cell, "span_width", 1) or 1
                                if sh > 1 or sw > 1:
                                    merge_count += 1
        return {
            "slide_count": slide_count,
            "table_count": table_count,
            "merge_count": merge_count,
            "max_rows": max_rows,
            "max_cols": max_cols,
            "total_cells": total_cells,
            "error": None,
        }
    except Exception as e:
        return {"error": f"{type(e).__name__}: {e}"}


def main() -> int:
    paths = find_pptx()
    print(f"총 {len(paths)}개 PPTX 발견, 분석 중...\n")

    rows: list[tuple] = []
    for p in paths:
        size = p.stat().st_size
        m = analyze(p)
        if m.get("error"):
            rows.append((p, size, 0, 0, 0, m["error"][:40]))
        else:
            rows.append((
                p,
                size,
                m["slide_count"],
                m["table_count"],
                m["merge_count"],
                f"{m['max_rows']}x{m['max_cols']}",
            ))

    # 표 있는 것 우선, 병합 셀 있는 것 더 우선, 크기 다양성 고려
    rows.sort(
        key=lambda r: (-r[4] if isinstance(r[4], int) else 0,
                       -r[3] if isinstance(r[3], int) else 0,
                       r[1]),
    )

    print(f"{'Size':>8} {'Slides':>6} {'Tables':>6} {'Merges':>6} {'MaxRC':>8}  Path")
    print("-" * 100)
    for p, size, slides, tables, merges, maxrc in rows[:30]:
        size_kb = f"{size//1024}K" if size < 1024*1024 else f"{size/(1024*1024):.1f}M"
        print(f"{size_kb:>8} {slides:>6} {tables:>6} {merges:>6} {str(maxrc):>8}  {p}")

    print(f"\n표 있는 파일: {sum(1 for r in rows if isinstance(r[3], int) and r[3] > 0)}")
    print(f"병합 셀 있는 파일: {sum(1 for r in rows if isinstance(r[4], int) and r[4] > 0)}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
