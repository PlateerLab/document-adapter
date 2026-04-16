"""합성 보고서 PPTX fixture — 분기 실적 보고서 레이아웃.

웹에서 공개 pptx 샘플 다운이 안정적이지 않아 (공유마당 JS-only, file-examples
HTML block 등) python-pptx 로 실전 분기 보고서 레이아웃을 합성.

레이아웃:
  slide 1: 표지 (제목, 작성자, 작성일, 부서)
  slide 2: 목차
  slide 3: 분기 KPI 요약 (표 1: 항목 | 목표 | 실적 | 달성률)
  slide 4: 부문별 실적 (표 2: 병합 헤더 + 세부 데이터)
  slide 5: 다음 분기 계획 (표 3: 과제 | 담당자 | 완료일)
"""
from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def make_quarterly_report(dst: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # --- slide 1: 표지 (양식 — fill_form 대상) ---
    s1 = prs.slides.add_slide(blank)
    title_box = s1.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "2026년 1분기 경영 실적 보고서"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True

    # 표지 양식 (라벨|값 4x2)
    t1 = s1.shapes.add_table(
        4, 2, Inches(3.5), Inches(4.5), Inches(6), Inches(2)
    ).table
    for i, (label, val) in enumerate([
        ("보고일자", ""),
        ("작성자", ""),
        ("작성부서", ""),
        ("승인자", ""),
    ]):
        t1.cell(i, 0).text = label
        t1.cell(i, 1).text = val

    # --- slide 2: 목차 ---
    s2 = prs.slides.add_slide(blank)
    box = s2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    box.text_frame.text = "목차"
    box.text_frame.paragraphs[0].font.size = Pt(32)
    toc = s2.shapes.add_table(
        5, 2, Inches(2), Inches(2), Inches(9), Inches(4)
    ).table
    for i, (num, item) in enumerate([
        ("1", "분기 KPI 요약"),
        ("2", "부문별 매출 실적"),
        ("3", "주요 성과 지표"),
        ("4", "다음 분기 계획"),
        ("5", "리스크 및 대응"),
    ]):
        toc.cell(i, 0).text = num
        toc.cell(i, 1).text = item

    # --- slide 3: 분기 KPI (표 — 양식성) ---
    s3 = prs.slides.add_slide(blank)
    box = s3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    box.text_frame.text = "1. 분기 KPI 요약"
    box.text_frame.paragraphs[0].font.size = Pt(28)

    kpi = s3.shapes.add_table(
        6, 4, Inches(0.5), Inches(1.5), Inches(12), Inches(4)
    ).table
    for i, hdr in enumerate(["항목", "목표", "실적", "달성률"]):
        kpi.cell(0, i).text = hdr
    kpi_rows = [
        ("매출", "", "", ""),
        ("영업이익", "", "", ""),
        ("신규고객", "", "", ""),
        ("고객만족도", "", "", ""),
        ("제품출시", "", "", ""),
    ]
    for r, row in enumerate(kpi_rows, start=1):
        for c, v in enumerate(row):
            kpi.cell(r, c).text = v

    # --- slide 4: 부문별 실적 (병합 있는 표) ---
    s4 = prs.slides.add_slide(blank)
    box = s4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    box.text_frame.text = "2. 부문별 매출 실적"
    box.text_frame.paragraphs[0].font.size = Pt(28)

    # 병합: row 0 은 연도 그룹 — (0,1) 2025년 (3병합), (0,4) 2026년 Q1 (2병합)
    perf = s4.shapes.add_table(
        6, 6, Inches(0.5), Inches(1.5), Inches(12), Inches(4)
    ).table
    perf.cell(0, 0).text = "부문"
    # 2025년 가로 3열 병합
    perf.cell(0, 1).merge(perf.cell(0, 3))
    perf.cell(0, 1).text = "2025년"
    # 2026 Q1 가로 2열 병합
    perf.cell(0, 4).merge(perf.cell(0, 5))
    perf.cell(0, 4).text = "2026년 Q1"

    perf.cell(1, 0).text = ""  # 부문 아래 빈칸
    for i, q in enumerate(["2Q", "3Q", "4Q", "목표", "실적"]):
        perf.cell(1, i + 1).text = q

    divisions = [
        ("사업A", "", "", "", "", ""),
        ("사업B", "", "", "", "", ""),
        ("사업C", "", "", "", "", ""),
        ("합계", "", "", "", "", ""),
    ]
    for r, row in enumerate(divisions, start=2):
        for c, v in enumerate(row):
            perf.cell(r, c).text = v

    # --- slide 5: 다음 분기 계획 ---
    s5 = prs.slides.add_slide(blank)
    box = s5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    box.text_frame.text = "3. 다음 분기 계획"
    box.text_frame.paragraphs[0].font.size = Pt(28)

    plan = s5.shapes.add_table(
        5, 4, Inches(0.5), Inches(1.5), Inches(12), Inches(4)
    ).table
    for i, hdr in enumerate(["우선순위", "과제", "담당자", "완료일"]):
        plan.cell(0, i).text = hdr
    for r in range(1, 5):
        for c in range(4):
            plan.cell(r, c).text = ""

    prs.save(str(dst))
    print(f"생성: {dst.name} ({dst.stat().st_size:,} bytes)")


if __name__ == "__main__":
    out = ROOT / "tests" / "fixtures" / "pptx" / "real" / "quarterly_report_synth.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    make_quarterly_report(out)
