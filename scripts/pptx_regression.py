"""PPTX round-trip regression harness.

tests/fixtures/pptx/real/ 하위의 모든 PPTX 파일에 대해:
  Stage A: bytes-identical copy (zipfile)
  Stage B: lxml XML 재직렬화 round-trip
  Stage C: adapter load → save (수정 없음) → 메트릭 보존
  Stage D: adapter 편집 → reload → 값/구조 보존

사용:
  python scripts/pptx_regression.py              # 모든 샘플 verify
  python scripts/pptx_regression.py --baseline   # 현재를 baseline으로 저장
  python scripts/pptx_regression.py --compare    # baseline 대비 drift 감지
"""
from __future__ import annotations

import argparse
import hashlib
import json
import re
import sys
import traceback
import warnings
import zipfile
from dataclasses import asdict, dataclass, field
from pathlib import Path

from lxml import etree

SCRIPT_DIR = Path(__file__).resolve().parent
ROOT = SCRIPT_DIR.parent
FIXTURES = ROOT / "tests" / "fixtures" / "pptx" / "real"
WORK = SCRIPT_DIR / "_poc_out" / "pptx_regression"
BASELINE = SCRIPT_DIR / "pptx_regression_baseline.json"


# ---------- 결과 모델 ----------

@dataclass
class Metrics:
    slide_count: int = 0
    table_count: int = 0
    merged_cell_count: int = 0
    total_rows: int = 0
    total_cols: int = 0
    total_cells: int = 0
    total_text_length: int = 0
    placeholder_count: int = 0
    error: str | None = None


@dataclass
class FixtureResult:
    name: str
    size_bytes: int
    sha256: str
    bytes_copy_ok: bool = False
    bytes_copy_identical: bool = False
    lxml_rt_ok: bool = False
    lxml_rt_reopen: bool = False
    python_pptx_open: bool = False
    adapter_rt_ok: bool = False
    adapter_rt_reopen: bool = False
    edit_ok: bool = False
    edit_reopen: bool = False
    edit_value_roundtrip: bool = False
    edit_structure_preserved: bool = False
    edit_target: tuple[int, int, int] | None = None
    edit_old_value: str | None = None
    edit_observed_value: str | None = None
    metrics_original: Metrics = field(default_factory=Metrics)
    metrics_bytes_copy: Metrics = field(default_factory=Metrics)
    metrics_lxml: Metrics = field(default_factory=Metrics)
    metrics_adapter: Metrics = field(default_factory=Metrics)
    metrics_edited: Metrics = field(default_factory=Metrics)
    errors: list[str] = field(default_factory=list)


# ---------- helpers ----------

def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()[:16]


# ---------- Stage A: bytes-identical copy ----------

def bytes_copy(src: Path, dst: Path) -> None:
    with zipfile.ZipFile(src, "r") as zin, zipfile.ZipFile(dst, "w") as zout:
        for info in zin.infolist():
            data = zin.read(info.filename)
            new_info = zipfile.ZipInfo(filename=info.filename, date_time=info.date_time)
            new_info.compress_type = info.compress_type
            new_info.external_attr = info.external_attr
            new_info.internal_attr = info.internal_attr
            new_info.create_system = info.create_system
            new_info.create_version = info.create_version
            new_info.extract_version = info.extract_version
            new_info.flag_bits = info.flag_bits
            zout.writestr(new_info, data)


def zips_identical(a: Path, b: Path) -> bool:
    with zipfile.ZipFile(a, "r") as za, zipfile.ZipFile(b, "r") as zb:
        names_a = za.namelist()
        names_b = zb.namelist()
        if names_a != names_b:
            return False
        for name in names_a:
            if za.read(name) != zb.read(name):
                return False
    return True


# ---------- Stage B: lxml round-trip ----------

def lxml_roundtrip(src: Path, dst: Path) -> None:
    parser = etree.XMLParser(remove_blank_text=False, strip_cdata=False)
    with zipfile.ZipFile(src, "r") as zin, zipfile.ZipFile(dst, "w") as zout:
        for info in zin.infolist():
            data = zin.read(info.filename)
            if info.filename.endswith((".xml", ".rels")):
                try:
                    tree = etree.fromstring(data, parser=parser)
                    data = etree.tostring(
                        tree,
                        xml_declaration=True,
                        encoding="UTF-8",
                        standalone=True,
                    )
                except etree.XMLSyntaxError:
                    pass
            new_info = zipfile.ZipInfo(filename=info.filename, date_time=info.date_time)
            new_info.compress_type = info.compress_type
            new_info.external_attr = info.external_attr
            zout.writestr(new_info, data)


# ---------- 메트릭 추출 ----------

PLACEHOLDER_RE = re.compile(r"\{\{\s*\w+\s*\}\}")


def extract_metrics(path: Path) -> Metrics:
    try:
        from pptx import Presentation
    except ImportError:
        return Metrics(error="python-pptx not installed")

    m = Metrics()
    try:
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            prs = Presentation(str(path))
        m.slide_count = len(prs.slides)

        text_buf: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                # 텍스트 수집
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text:
                                text_buf.append(run.text)
                # 표
                if shape.has_table:
                    m.table_count += 1
                    tbl = shape.table
                    rows = len(tbl.rows)
                    cols = len(tbl.columns)
                    m.total_rows += rows
                    m.total_cols += cols
                    for row in tbl.rows:
                        for cell in row.cells:
                            m.total_cells += 1
                            if getattr(cell, "is_merge_origin", False):
                                sh = getattr(cell, "span_height", 1) or 1
                                sw = getattr(cell, "span_width", 1) or 1
                                if sh > 1 or sw > 1:
                                    m.merged_cell_count += 1
                            # 셀 내부 텍스트도 포함
                            if cell.text_frame:
                                for p in cell.text_frame.paragraphs:
                                    for r in p.runs:
                                        if r.text:
                                            text_buf.append(r.text)

        full_text = "\n".join(text_buf)
        m.total_text_length = len(full_text)
        m.placeholder_count = len(PLACEHOLDER_RE.findall(full_text))
    except Exception as e:
        m.error = f"{type(e).__name__}: {e}"
    return m


# ---------- 메인 검증 ----------

def verify_fixture(path: Path) -> FixtureResult:
    result = FixtureResult(
        name=path.name,
        size_bytes=path.stat().st_size,
        sha256=sha256_file(path),
    )

    WORK.mkdir(parents=True, exist_ok=True)
    copy_dst = WORK / f"copy_{path.name}"
    lxml_dst = WORK / f"lxml_{path.name}"

    result.metrics_original = extract_metrics(path)
    if result.metrics_original.error is None:
        result.python_pptx_open = True

    # Stage A
    try:
        bytes_copy(path, copy_dst)
        result.bytes_copy_ok = True
        result.bytes_copy_identical = zips_identical(path, copy_dst)
        result.metrics_bytes_copy = extract_metrics(copy_dst)
    except Exception as e:
        result.errors.append(f"bytes_copy: {type(e).__name__}: {e}")

    # Stage B
    try:
        lxml_roundtrip(path, lxml_dst)
        result.lxml_rt_ok = True
        result.metrics_lxml = extract_metrics(lxml_dst)
        result.lxml_rt_reopen = result.metrics_lxml.error is None
    except Exception as e:
        result.errors.append(f"lxml_rt: {type(e).__name__}: {e}")

    # Stage C
    adapter_dst = WORK / f"adapter_{path.name}"
    try:
        from document_adapter import load

        adapter = load(path)
        try:
            adapter.save(adapter_dst)
        finally:
            adapter.close()
        result.adapter_rt_ok = True
        result.metrics_adapter = extract_metrics(adapter_dst)
        result.adapter_rt_reopen = result.metrics_adapter.error is None
    except Exception as e:
        result.errors.append(f"adapter_rt: {type(e).__name__}: {e}\n{traceback.format_exc()[:300]}")

    # Stage D
    edit_dst = WORK / f"edit_{path.name}"
    try:
        from document_adapter import load

        sentinel = "__PPTX_ADAPTER_EDIT_CHECK_2026_04_16__"

        adapter = load(path)
        try:
            tables = adapter.get_tables(preview_rows=1000)
            target = None
            for t in tables:
                for r, row in enumerate(t.preview):
                    for c, val in enumerate(row):
                        if val is not None:
                            target = (t.index, r, c)
                            break
                    if target:
                        break
                if target:
                    break

            if target is None:
                result.errors.append("edit: 편집 가능한 anchor 셀 없음")
                return result

            tidx, r, c = target
            result.edit_target = target
            result.edit_old_value = adapter.set_cell(tidx, r, c, sentinel)
            adapter.save(edit_dst)
            result.edit_ok = True
        finally:
            adapter.close()

        adapter2 = load(edit_dst)
        try:
            cell = adapter2.get_cell(tidx, r, c)
            result.edit_reopen = True
            result.edit_observed_value = cell.text.strip()
            result.edit_value_roundtrip = (cell.text.strip() == sentinel)
        finally:
            adapter2.close()

        result.metrics_edited = extract_metrics(edit_dst)
        m_orig = result.metrics_original
        m_edit = result.metrics_edited
        if m_edit.error is None:
            result.edit_structure_preserved = (
                m_orig.slide_count == m_edit.slide_count
                and m_orig.table_count == m_edit.table_count
                and m_orig.merged_cell_count == m_edit.merged_cell_count
                and m_orig.placeholder_count == m_edit.placeholder_count
                and m_orig.total_rows == m_edit.total_rows
                and m_orig.total_cols == m_edit.total_cols
                and m_orig.total_cells == m_edit.total_cells
            )
    except Exception as e:
        result.errors.append(f"edit: {type(e).__name__}: {e}\n{traceback.format_exc()[:300]}")

    return result


def metrics_match(a: Metrics, b: Metrics) -> bool:
    if a.error or b.error:
        return False
    return (
        a.slide_count == b.slide_count
        and a.table_count == b.table_count
        and a.merged_cell_count == b.merged_cell_count
        and a.total_rows == b.total_rows
        and a.total_cols == b.total_cols
        and a.total_cells == b.total_cells
        and a.total_text_length == b.total_text_length
        and a.placeholder_count == b.placeholder_count
    )


def print_result(r: FixtureResult) -> None:
    all_ok = (
        r.bytes_copy_identical
        and r.lxml_rt_reopen
        and r.adapter_rt_reopen
        and r.edit_value_roundtrip
        and r.edit_structure_preserved
        and metrics_match(r.metrics_original, r.metrics_lxml)
        and metrics_match(r.metrics_original, r.metrics_adapter)
    )
    status = "✅" if all_ok else ("⚠️" if r.python_pptx_open else "❌")
    print(f"\n{status} {r.name} ({r.size_bytes:,} bytes, sha={r.sha256})")
    m = r.metrics_original
    if m.error:
        print(f"    원본 read 실패: {m.error}")
        return
    print(
        f"    원본: slides={m.slide_count}, tables={m.table_count}, "
        f"merges={m.merged_cell_count}, cells={m.total_cells}, "
        f"text_len={m.total_text_length:,}, placeholders={m.placeholder_count}"
    )
    print(
        f"    Stage A (bytes copy)  : ok={r.bytes_copy_ok}, "
        f"identical={r.bytes_copy_identical}, "
        f"metrics_match={metrics_match(m, r.metrics_bytes_copy)}"
    )
    print(
        f"    Stage B (lxml rt)     : ok={r.lxml_rt_ok}, "
        f"reopen={r.lxml_rt_reopen}, "
        f"metrics_match={metrics_match(m, r.metrics_lxml)}"
    )
    print(
        f"    Stage C (adapter rt)  : ok={r.adapter_rt_ok}, "
        f"reopen={r.adapter_rt_reopen}, "
        f"metrics_match={metrics_match(m, r.metrics_adapter)}"
    )
    if r.edit_target:
        t, row, col = r.edit_target
        print(
            f"    Stage D (adapter edit): ok={r.edit_ok}, "
            f"reopen={r.edit_reopen}, "
            f"value_roundtrip={r.edit_value_roundtrip}, "
            f"structure_preserved={r.edit_structure_preserved}"
        )
        old_disp = (r.edit_old_value or "")[:60]
        print(
            f"      target=T{t}({row},{col}), old={old_disp!r}, "
            f"observed={r.edit_observed_value!r}"
        )
    else:
        print(f"    Stage D (adapter edit): SKIPPED (no anchor cell)")
    if r.errors:
        for e in r.errors:
            print(f"    ERR: {e}")


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--baseline", action="store_true")
    ap.add_argument("--compare", action="store_true")
    ap.add_argument("--json", action="store_true")
    args = ap.parse_args()

    if not FIXTURES.exists():
        print(f"fixtures 디렉토리 없음: {FIXTURES}")
        return 2

    fixtures = sorted(FIXTURES.glob("*.pptx"))
    if not fixtures:
        print(f"PPTX 파일 없음: {FIXTURES}")
        return 2

    print(f"{len(fixtures)}개 fixture 검증 중...")
    results: list[FixtureResult] = []
    for p in fixtures:
        r = verify_fixture(p)
        results.append(r)
        print_result(r)

    print("\n" + "=" * 70)
    total = len(results)
    opened = sum(1 for r in results if r.python_pptx_open)
    bytes_ok = sum(1 for r in results if r.bytes_copy_identical)
    lxml_ok = sum(
        1 for r in results
        if r.lxml_rt_reopen and metrics_match(r.metrics_original, r.metrics_lxml)
    )
    adapter_ok = sum(
        1 for r in results
        if r.adapter_rt_reopen and metrics_match(r.metrics_original, r.metrics_adapter)
    )
    edit_ok = sum(
        1 for r in results
        if r.edit_value_roundtrip and r.edit_structure_preserved
    )
    print(f"요약: {total} fixtures")
    print(f"  python-pptx로 원본 열림        : {opened}/{total}")
    print(f"  bytes copy 100% identical      : {bytes_ok}/{total}")
    print(f"  lxml round-trip 메트릭 보존    : {lxml_ok}/{total}")
    print(f"  adapter round-trip 메트릭 보존 : {adapter_ok}/{total}")
    print(f"  adapter 편집-후 값/구조 보존   : {edit_ok}/{total}")

    report = {
        "fixtures": [asdict(r) for r in results],
        "summary": {
            "total": total, "opened": opened, "bytes_ok": bytes_ok,
            "lxml_ok": lxml_ok, "adapter_ok": adapter_ok, "edit_ok": edit_ok,
        },
    }

    if args.baseline:
        BASELINE.write_text(json.dumps(report, ensure_ascii=False, indent=2))
        print(f"\nbaseline 저장: {BASELINE}")

    if args.compare:
        if not BASELINE.exists():
            print(f"\nbaseline 없음")
            return 3
        old = json.loads(BASELINE.read_text())
        old_by_name = {f["name"]: f for f in old["fixtures"]}
        drift = []
        for r in results:
            old_r = old_by_name.get(r.name)
            if old_r is None:
                drift.append(f"신규: {r.name}")
                continue
            for key in ("bytes_copy_identical", "lxml_rt_reopen",
                        "adapter_rt_reopen", "edit_value_roundtrip",
                        "edit_structure_preserved"):
                if getattr(r, key) != old_r.get(key):
                    drift.append(f"{r.name}: {key} {old_r.get(key)} → {getattr(r, key)}")
        if drift:
            print(f"\n⚠️ drift ({len(drift)}):")
            for d in drift:
                print(f"  - {d}")
            return 1
        else:
            print("\n✅ baseline과 동일")

    if args.json:
        print("\n" + json.dumps(report, ensure_ascii=False, indent=2))

    return 0 if (adapter_ok == total and edit_ok == total) else 1


if __name__ == "__main__":
    sys.exit(main())
